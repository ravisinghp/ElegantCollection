import aiohttp
from datetime import datetime
from typing import List,Optional
from app.db.repositories.sharepoint_repo import SharepointRepo
import os,base64
import hashlib
from dotenv import load_dotenv
import logging
import re, json, io, PyPDF2, docx
from pptx import Presentation
from rapidfuzz import fuzz
from openai import OpenAI
from decimal import Decimal
from datetime import date, datetime,timedelta,timezone
from fastapi import APIRouter, HTTPException,Query,Request
import pandas as pd
import asyncio
from app.services.usersmailservice import FIELDS_TO_COMPARE
from app.utils.image_ocr import extract_text_from_image_bytes
from collections import defaultdict
from app.db.repositories.sync_client_po_repo import MSSQLRepo


# Load env
load_dotenv()
GRAPH_API = os.getenv("GRAPH_API")

# -------------- SharePoint Config start-------------- #
SHAREPOINT_SITE_URL = os.getenv("SHAREPOINT_SITE_URL")
SHAREPOINT_SITE_PATH = os.getenv("SHAREPOINT_SITE_PATH")
LIBRARY_NAME = os.getenv("LIBRARY_NAME")
# ------------- SharePoint Config end----------------- #

#---------------OpenAI Client------------------
openai_client = OpenAI(api_key=os.getenv("OPENAI_API_KEY"))
#----------OpenAI Client end ------------------

logger = logging.getLogger("sharepoint_service")
logger.setLevel(logging.INFO)


class SharepointService:
    def __init__(self, sp_repo: SharepointRepo):
            self.sp_repo = sp_repo
            
    @staticmethod
    def _generate_po_pdf(df, filename_prefix):
        from reportlab.lib.pagesizes import A4
        from reportlab.pdfgen import canvas
        import io

        output = io.BytesIO()
        pdf = canvas.Canvas(output, pagesize=A4)

        width, height = A4
        x_start = 30
        y = height - 40

        # Header
        for col in df.columns:
            pdf.drawString(x_start, y, str(col))
            x_start += 100

        y -= 20
        x_start = 30

        # Rows
        for _, row in df.iterrows():
            for val in row:
                pdf.drawString(x_start, y, str(val))
                x_start += 100

            y -= 20
            x_start = 30

            if y < 40:
                pdf.showPage()
                y = height - 40

        pdf.save()
        output.seek(0)

        return (
            output,
            f"{filename_prefix}.pdf",
            "application/pdf"
        )
        
        
    # ---------------- GET SITES BY USER EMAIL ---------------- #
    async def get_sites_by_user_email(self, access_token: str, user_email: str):
        headers = {"Authorization": f"Bearer {access_token}"}
        sites = []
        seen_site_ids = set()   #add this
 
        async with aiohttp.ClientSession() as session:
 
            #Group memberships
            async with session.get(
                "https://graph.microsoft.com/v1.0/me/memberOf",
                headers=headers
            ) as resp:
                data = await resp.json()
 
                for item in data.get("value", []):
                    if item["@odata.type"] == "#microsoft.graph.group":
                        group_id = item["id"]
 
                        async with session.get(
                            f"https://graph.microsoft.com/v1.0/groups/{group_id}/sites/root",
                            headers=headers
                        ) as site_resp:
                            if site_resp.status == 200:
                                site = await site_resp.json()
 
                                if site["id"] not in seen_site_ids: #prevent duplicate
                                    seen_site_ids.add(site["id"])
                                    sites.append({
                                        "id": site["id"],
                                        "name": site["displayName"],
                                        "webUrl": site["webUrl"]
                                    })
 
            #Followed sites
            async with session.get(
                "https://graph.microsoft.com/v1.0/me/followedSites",
                headers=headers
            ) as resp:
                if resp.status == 200:
                    data = await resp.json()
                    for site in data.get("value", []):
 
                        if site["id"] not in seen_site_ids:  #prevent duplicate
                            seen_site_ids.add(site["id"])
                            sites.append({
                                "id": site["id"],
                                "name": site["displayName"],
                                "webUrl": site["webUrl"]
                            })
 
        logger.info(f"Found {len(sites)} unique sites for user {user_email}")
        return sites
 
           
    #Fetching Total Numbers of Attachments on User Dashboard
    async def get_documents_analyzed_by_user_id(user_id: int, request: Request):
        try:
            return await SharepointRepo.fetch_documents_analyzed_by_user_id(user_id,  request)
        except Exception as e:
            return None
 

            
    #Fetching Total Numbers of Attachments on User Dashboard
    async def get_documents_analyzed_by_user_id(user_id: int, request: Request):
        try:
            return await SharepointRepo.fetch_documents_analyzed_by_user_id(user_id,  request)
        except Exception as e:
            return None
        
    
    async def get_drive_id(self, access_token: str, site_id: str):
        headers = {"Authorization": f"Bearer {access_token}"}

        async with aiohttp.ClientSession() as session:
            async with session.get(
                f"{GRAPH_API}/sites/{site_id}/drives",
                headers=headers
            ) as resp:

                if resp.status != 200:
                    text = await resp.text()
                    logger.error(f"Failed to get drives: {resp.status} | {text}")
                    raise Exception("Failed to get drives")

                data = await resp.json()

                for drive in data.get("value", []):
                    if drive.get("name") == LIBRARY_NAME:
                        return drive.get("id")

        raise Exception(f"Library '{LIBRARY_NAME}' not found")


    # ---------------- LIST ALL FOLDERS RECURSIVELY ---------------- #
    async def list_folders_recursive(self, access_token: str, drive_id: str, folder_path: str = ""):
        folders = []

        async def fetch_children(path: str, parent_path: str = ""):
            url = f"{GRAPH_API}/drives/{drive_id}/root"
            if path:
                url += f":/{path}:/children"
            else:
                url += "/children"

            headers = {"Authorization": f"Bearer {access_token}"}
            async with aiohttp.ClientSession() as session:
                async with session.get(url, headers=headers) as resp:
                    if resp.status != 200:
                        text = await resp.text()
                        logger.error(f"Failed to fetch folders: {resp.status} | {text}")
                        return
                    data = await resp.json()
                    for item in data.get("value", []):
                        if "folder" in item:
                            folder_name = item.get("name")
                            full_path = f"{parent_path}/{folder_name}" if parent_path else folder_name
                            folders.append({"id": item.get("id"), "name": folder_name, "path": full_path})
                            await fetch_children(full_path, full_path)

        await fetch_children(folder_path)
        return folders


    # ---------------- FETCH DRIVE FILES ---------------- #
    async def fetch_drive_files(
        self,
        access_token: str,
        drive_id: str,
        folder_path: str = "",
        from_date: str = None,
        to_date: str = None,
    ) -> List[dict]:

        headers = {"Authorization": f"Bearer {access_token}"}

        async def walk(path: str):
            if folder_path:
                url = f"{GRAPH_API}/drives/{drive_id}/root:/{folder_path}:/children"
            else:
                url = f"{GRAPH_API}/drives/{drive_id}/root/children"

            async with aiohttp.ClientSession(headers=headers) as session:
                async with session.get(url) as resp:
                    if resp.status != 200:
                        text = await resp.text()
                        logger.error(f"Failed: {resp.status} | {text}")
                        return []

                    data = await resp.json()

            collected_files = []

            for item in data.get("value", []):
                if "file" in item:
                    collected_files.append(item)

            return collected_files


        result = await walk(folder_path)

        for f in result:
            print(f)
        
        

        # ---------- DATE FILTER ----------
        if from_date or to_date:
            # from_dt = datetime.fromisoformat(from_date) if from_date else None
            # to_dt = datetime.fromisoformat(to_date) if to_date else None
            from_dt = (
                datetime.fromisoformat(from_date).replace(tzinfo=timezone.utc)
                if from_date else None
            )

            to_dt = (
                datetime.fromisoformat(to_date)
                .replace(tzinfo=timezone.utc)
                + timedelta(days=1)
                if to_date else None
            )

            filtered = []
            for f in result:
                created_str = f.get("createdDateTime")

                if not created_str:
                    logger.warning("File missing createdDateTime: %s", f.get("id"))
                    continue

                try:
                    created_dt = datetime.fromisoformat(
                        created_str.replace("Z", "+00:00")
                    )
                except Exception as e:
                    logger.warning(
                        "Invalid createdDateTime (%s) for file %s",
                        created_str,
                        f.get("id"),
                    )
                    continue

                if from_dt and created_dt < from_dt:
                    continue
                if to_dt and created_dt >= to_dt:  # >= is intentional
                    continue

                filtered.append(f)

            result = filtered

        return result

    # ---------------- UTILS ---------------- #
    @staticmethod
    def graph_datetime_to_mysql(dt_str: str | None) -> str | None:
        if not dt_str:
            return None
        return datetime.fromisoformat(
            dt_str.replace("Z", "+00:00")
        ).strftime("%Y-%m-%d %H:%M:%S")

    # ---------------- FILE HASHING ---------------- #
    @staticmethod
    async def generate_file_hash(file_bytes: bytes) -> str:
        return hashlib.sha256(file_bytes).hexdigest()
    
    
    
    @staticmethod
    def strip_item_column_noise(desc: str) -> str:
        if not desc:
            return desc

        BAD_PREFIXES = [
            "expected delivery",
            "delivery date",
            "qty",
            "quantity",
            "item description",
            "description",
        ]

        d = desc.strip()

        for p in BAD_PREFIXES:
            if d.lower().startswith(p):
                d = d[len(p):].strip()

        return d

    # ---------------- TEXT EXTRACTION ---------------- #
    def xlrd_cell_to_str(cell, workbook) -> str:
        import xlrd
        from datetime import datetime

        try:
            # Proper date handling
            if cell.ctype == xlrd.XL_CELL_DATE:
                dt = xlrd.xldate_as_datetime(cell.value, workbook.datemode)
                return dt.strftime("%Y-%m-%d")

            # Handle numbers that might actually be dates
            if cell.ctype == xlrd.XL_CELL_NUMBER:
                # Try converting to date if possible
                try:
                    dt = xlrd.xldate_as_datetime(cell.value, workbook.datemode)
                    return dt.strftime("%Y-%m-%d")
                except Exception:
                    return str(cell.value)

            if cell.ctype == xlrd.XL_CELL_EMPTY:
                return ""

            return str(cell.value)

        except Exception:
            return str(cell.value)

    @staticmethod
    async def extract_text_from_bytes(content_bytes: bytes, filename: str, content_type: str) -> str | None:
        """
        Fast, async-safe text extraction from attachments.
        Runs heavy parsing in a background thread.
        """
        ext = (filename or "").lower()
        ct = (content_type or "").lower()

        def parse_attachment():
            try:
                # Text files
                if ct.startswith("text/") or ext.endswith((".txt", ".md", ".csv", ".log")):
                    return content_bytes.decode("utf-8", errors="ignore")

                # PDF files
                elif ct == "application/pdf" or ext.endswith(".pdf"):
                    reader = PyPDF2.PdfReader(io.BytesIO(content_bytes))
                    return " ".join((p.extract_text() or "") for p in reader.pages)

                # Word documents
                elif ct in (
                    "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                    "application/msword",
                ) or ext.endswith((".docx", ".doc")):
                    doc = docx.Document(io.BytesIO(content_bytes))
                    if ct in (
                        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                        "application/msword",
                    ) or ext.endswith((".docx", ".doc")):

                        doc = docx.Document(io.BytesIO(content_bytes))

                        text_parts = []

                        # ---- paragraphs ----
                        for p in doc.paragraphs:
                            if p.text.strip():
                                text_parts.append(p.text.strip())

                        # ---- tables (CRITICAL FIX) ----
                        for table in doc.tables:
                            for row in table.rows:
                                cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                                if cells:
                                    text_parts.append(" | ".join(cells))

                        return "\n".join(text_parts)

                # PowerPoint files
                elif ct in ("application/vnd.openxmlformats-officedocument.presentationml.presentation",
                            "application/vnd.ms-powerpoint") or ext.endswith((".pptx", ".ppt")):
                    prs = Presentation(io.BytesIO(content_bytes))
                    return " ".join(
                        shape.text
                        for slide in prs.slides
                        for shape in slide.shapes
                        if hasattr(shape, "text")
                    )

                elif ct in (
                    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                    "application/vnd.ms-excel"
                ) or ext.endswith((".xlsx", ".xls")):
                    all_text = []
                    if ext.endswith(".xlsx"):
                        from openpyxl import load_workbook
                        from datetime import datetime

                        wb = load_workbook(io.BytesIO(content_bytes), data_only=True)

                        for sheet in wb.worksheets:   # FIX
                            all_text.append(f"Sheet: {sheet.title}")

                            for row in sheet.iter_rows():
                                values = []
                                for cell in row:
                                    val = cell.value

                                    # Proper date handling
                                    if isinstance(val, datetime):
                                        val = val.strftime("%Y-%m-%d")

                                    values.append(str(val) if val is not None else "")

                                row_text = " | ".join(v for v in values if v.strip())
                                if row_text:
                                    all_text.append(row_text)
                    elif ext.endswith(".xls"):
                        import xlrd
                        wb = xlrd.open_workbook(file_contents=content_bytes)
                        for sheet in wb.sheets():
                            all_text.append(f"Sheet: {sheet.name}")
                            for row_idx in range(sheet.nrows):
                                row_vals = [
                                    SharepointService.xlrd_cell_to_str(sheet.cell(row_idx, col_idx), wb)
                                    for col_idx in range(sheet.ncols)
                                ]
                                row_text = " | ".join(v for v in row_vals if v.strip())
                                if row_text.strip():
                                    all_text.append(row_text)
                    return "\n".join(all_text)
                
                # Other formats (images/ocr) - optional
                elif filename.endswith(SharepointService.IMAGE_EXTENSIONS):
                # IMAGE → OCR
                    try:
                        return extract_text_from_image_bytes(content_bytes)
                    except Exception as e:
                        logger.error(f"OCR failed for {filename}: {e}")
                        return ""
                else:
                    return None                     

            except Exception:
                return None
            
        return await asyncio.to_thread(parse_attachment)

    # ---------------- KEYWORD DETECTION ---------------- #
    @staticmethod
    def normalize_keyword(k: str) -> str:
        return re.sub(r"\s+", " ", k.strip().lower())
    
    
    # ---------------- KEYWORD DETECTION ---------------- #
    KEYWORD_REGEX_MAP = {
        "po_number": [
            r"\bpo\s*(no|number|#|id)\b",
            r"\bp\.o\.\s*(no|number|#)?\b",
            r"\bpurchase\s*order\s*(no|number)?\b",
            r"\b(po)\s*(?:no|number|#|id)?\s*\d+\b",
            r"\b(po)\d+\b",
            r"\b(p\.o\.)\s*(?:no|number|#)?\s*\d*\b",
            r"\b(purchase\s*order)\s*(?:no|number)?\s*\d*\b",
            r"\b(order)\s*(?:no|number|#|id)?\s*\d+\b",
            r"\b(order|orders|ordered|reorder|re-?order|ordering|purchase|purchases|po|attached|attachment|attachments|enclosed)\b"
        ],

        "customer_name": [
            r"\bcustomer\s*name\b",
            r"\bbuyer\b",
            r"\bclient\b",
            r"\bparty\s*name\b"
        ],

        "vendor_number": [
            r"\bvendor\s*(no|number|code)\b",
            r"\bsupplier\s*(no|number|code)\b"
        ],

        "po_date": [
            r"\bpo\s*date\b",
            r"\border\s*date\b",
            r"\bdate\s*of\s*order\b"
        ],

        "delivery_date": [
            r"\bdelivery\s*date\b",
            r"\bexpected\s*delivery\b",
            r"\bdispatch\s*date\b"
        ],

        "cancel_date": [
            r"\bcancel\s*date\b",
            r"\bexpiry\s*date\b"
        ],

        "gold_karat": [
            r"\b\d{2}\s*K\b",
            r"\b\d{2}\s*KT\b",
            r"\bAG\s*\d{3}\b",
            r"\bgold\s*karat\b",
            r"\bpurity\b"
        ],

        "ec_style_number": [
            r"\bec\s*style\s*(no|number)\b",
            r"\bec\s*style\b"
        ],

        "customer_style_number": [
            r"\bcustomer\s*style\s*(no|number)\b",
            r"\bdesign\s*(no|number)\b"
        ],

        "color": [
            r"\b(yellow|white|rose)\s*gold\b",
            r"\bcolor\b"
        ],

        "quantity": [
            r"\bqty\b",
            r"\bquantity\b",
            r"\bpcs\b",
            r"\bnos\b",
            r"\bpieces\b"
        ],

        "description": [
            r"\bdescription\b",
            r"\bitem\s*details\b",
            r"\bproduct\s*details\b"
        ],

        "gold_lock": [
            r"\bgold\s*lock\b",
            r"\blocked\s*gold\b",
            r"\block\s*status\b"
        ],
    }

    async def detect_keywords(self, text: str, db_keywords: list[str]):
        if not text or not text.strip():
            return [], None

        text_l = text.lower()
        detected = set()

        # ---------------- 1. DB keywords (exact-ish) ----------------
        for kw in db_keywords:
            if self.normalize_keyword(kw) in text_l:
                detected.add(kw)

        # ---------------- 2. Jewellery-aware regex detection ----------------
        for field, patterns in self.KEYWORD_REGEX_MAP.items():
            for pat in patterns:
                if re.search(pat, text_l, re.IGNORECASE):
                    detected.add(field)
                    break

        if detected:
            return sorted(detected), "REGEX_MATCH"

        return [], None
    

    # ---------------- PO EXTRACTION ---------------- #
    PO_FIELD_NAMES = [
        "po_number",
        "customer_name",
        "vendor_number",
        "po_date",
        "delivery_date",
        "cancel_date",
        "ec_style_number",
        "customer_style_number",
        "quantity",
        "gold_karat",
        "color",
        "description",
        "gold_lock"
    ]

    EMPTY_PO = {k: None for k in PO_FIELD_NAMES}
    
    #--------------------Regex-----------------------------
    PO_REGEX_PATTERNS = {

        # ---------------- PO NUMBER ----------------
        "po_number": [
            r"(?:po_number|po_no)\s*:\s*(PO[\w\-_/]+)",
            r"(?:po\s*number|po\s*no|po#|po\s*#|p\.o\.|purchase\s*order|po)\s*[:\-]?\s*(PO[\w\-_/]+)",
            r"\b(PO[\s\-_:]*[0-9]{1,}[A-Z0-9\/_.\-]*)",
            r"(?:po\s*number|po\s*no|po#|p\.o\.|purchase\s*order)\s*[:\-]?\s*(PO[\- ]?[A-Z0-9\/_.\-]+)",
            r"(?:po\s*number|po\s*no|po#|p\.o\.|purchase\s*order)\s*[:\-]?\s*([A-Z]{1,5}-\d{4,}-\d+)",
            r"(?:po_number|po_no|po\s*number|po#|p\.o\.)\s*[:#]?\s*\n?\s*([A-Z0-9\-_/]+)",
        ],

        # ---------------- CUSTOMER NAME ----------------
        "customer_name": [
            r"(?i)\bcustomer\s*(?:name)?\b\s*\|\s*([A-Za-z0-9&.,\-/ ]{2,})",
            r"(?i)\bcustomer\s*(?:name)?\s*[:\-]\s*\n\s*([A-Za-z0-9&.,\-/ ]{2,})",
            r"(?i)\bcustomer\s*(?:name)?\s*[:\-]\s*([A-Za-z0-9&.,\-/ ]{2,})",
            r"(?i)\bship\s*to\s*:\s*\n\s*([A-Za-z0-9&.,\-/ ]{2,})",
            r"(?i)(?:ship\s*to|bill\s*to|deliver\s*to)\s*[:\-]\s*([A-Za-z0-9&.,\-/ ]{2,})",
            r"(?i)(?:customer|buyer|client)\s*[:\-]?\s*([A-Za-z0-9&.,\-/ ]+?)(?=\s+(?:vendor|po|date|delivery|qty|quantity|gold|color|description)\b|$)",
        ],

        # ---------------- VENDOR NUMBER ----------------
        "vendor_number": [
            r"Vendor\s*ID[\s\S]{0,80}\b(V\d{4,})\b",
            r"(?:vendor[_\s]*(?:number|no|id)|supplier[_\s]*(?:number|no|code))\s*[:\-#]?\s*([A-Za-z0-9\-_./]+)",
            r"(?:vendor\s*id|vendor\s*number)\s*[:\-]?\s*\n\s*([A-Za-z0-9\-_./]+)",
            r"\bvendor\b\s*[:\-]?\s*([A-Za-z0-9\-_./]+)",
            r"\b(?:Vendor\s*ID|VNo|V-ID)\s*[:#\-\s]?\s*([A-Za-z0-9\-_./]+)",
            r"(?:supplier\s*(?:no|number|code))\s*[:\-]?\s*([A-Za-z0-9\-_./]+)",
        ],                                                              

        # ---------------- PO DATE ----------------
        "po_date": [
            r"(?:po\s*date|order\s*date|date)\s*[:\-]?\s*(\d{4}-\d{1,2}-\d{1,2})",
            r"po_date\s*:\s*(\d{4}-\d{2}-\d{2})",
            r"date\s*[:\-]?\s*(\d{4}-\d{2}-\d{2})",
            r"(?:Purchase\s+Order\s+Date|P\.O\.\s*Date)\s*[:\-]?\s*(\d{1,2}/\d{1,2}/\d{2})",
            r"P\.?\s*O\.?\s*Date\s*[:\-]?\s*(\d{1,2}-[A-Za-z]{3}-\d{4})",
            r"(?:purchase\s*order\s*date)\s*\n\s*(\d{1,2}/\d{1,2}/\d{2})",
            r"(?:purchase\s*order\s*date)\s*[:\-]?\s*(\d{1,2}-[A-Za-z]{3}-\d{4})",
            r"(?:po\s*date|order\s*date|date)\s*[:\-]?\s*([0-9]{1,2}-[A-Za-z]{3}-[0-9]{4})",
            r"P\.O\.\s+Date\s*:\s*(\d{2}-[A-Za-z]{3}-\d{4})",
            r"P\.O\.\s*Date\s*[:\-]?\s*(\d{1,2}-[A-Za-z]{3}-\d{4})",
            r"(?:po\s*date|order\s*date|date)\s*[:\-]?\s*(\d{4}-\d{1,2}-\d{1,2})",
            r"\b(\d{1,2}/\d{1,2}/\d{2})\b",
        ],

        #---------------- DELIVERY DATE ----------------
        "delivery_date": [
            r"\b(?:DELIVERY\s*DATE|DUE\s*DATE|delivery_date)\b[\s\S]{0,100}?[:\s]*([\dA-Za-z/.-]{4,20})",
            r"(?:delivery\s*date|expected\s*delivery|delivery_date|due\s*date)\s*[:\-]?\s*(\d{4}-\d{2}-\d{2})",
            r"(?:delivery\s*date|expected\s*delivery|delivery_date|due\s*date)\s*[:\-]?\s*(\d{2}-[A-Za-z]{3}-\d{4})",
        ], 
    
        # ---------------- CANCEL DATE ----------------
        "cancel_date": [
            r"(?:cancel\s*date|cancellation\s*date)\s*[:\-]?\s*(\d{4}-\d{1,2}-\d{1,2})",
            r"cancel_date\s*:\s*(\d{4}-\d{2}-\d{2})",
        ],
    
        # ---------------- EC STYLE NUMBER ----------------
        "ec_style_number": [
        r"(?:ec\s*style\s*number|ec\s*style|ec\s*no)\s*[:\-]?\s*([A-Z0-9.\-]+)",
        r"(?:ec_style_number|ec_style_no)\s*:\s*([A-Z0-9.\-]+)",
        ],
    
        # ---------------- CUSTOMER STYLE NUMBER ----------------
        "customer_style_number": [
            r"(?:customer\s*style\s*number|customer\s*style|cust\s*style)\s*[:\-]?\s*([A-Z0-9\-]+)",
            r"(?:customer_style_number|customer_style_no)\s*:\s*([A-Z0-9\-]+)",
        ],
    
        # ---------------- QUANTITY ----------------
        "quantity": [
            r"\bQUANTITY\b[\s\S]{0,100}\n\s*([A-Za-z0-9 ,\-–\.]{10,})",
            r"(?i)\bquantity\b\s*[\r\n]+\s*(\d+)\b",
            r"\n\s*(\d+)\s+(?:EA|PCS|PC)\b",
            r"(?i)(?:qty|quantity|pcs|pieces)\s*[:\-]?\s*(\d+)",
            r"(?i)\b(\d+)\s*(?:pcs|pieces|nos)\b",
        ],
    
        # ---------------- GOLD KARAT ----------------
        "gold_karat": [
            r"(?i)\bMETAL\b[\s\S]{0,60}?\b([^\s|,\n\r]+)",
            r"(?i)(?:gold\s*karat|gold_karat|gold\s*purity|karat|metal|kt)\s*[:\-]?\s*([^\n\r]+?)(?=\s+(?:ec\s*style|customer\s*style|quantity|color|description|gold\s*lock)\s*[:\-]|$)",
            r"(?i)(?:gold\s*karat|gold_karat|karat|metal)\s*[\r\n]+\s*([^\r\n]+)",
        ],              
    
        # ---------------- COLOR ----------------
        "color": [
            r"(?:color|colour)\s*[:\-]?\s*([A-Za-z0-9\s\-]+)",
            r"(?:color|colour)\s*[:\-]?\s*([A-Za-z0-9\s\-]+?)(?=\s+(?:quantity|gold|karat|description|remarks|details)\s*:|$)",
            r"color\s*[:\-]?\s*([A-Za-z0-9\s\-]+)"
        ],

        #---------------- DESCRIPTION ----------------
        "description": [
            r"\bDESCRIPTION\b[\s\S]{0,100}\n\s*([A-Za-z0-9 ,\-–\.]{10,})",
            r"(?:item\s*description|description)\s*[:\-]?\s*([A-Za-z][A-Za-z\s\-–]+)",
            r"[A-Z0-9\-]+\s+\d+KW\s+([A-Za-z ].*?SIZE:\s*[0-9.]+)",
        ],

        # ---------------- GOLD LOCK ----------------
        "gold_lock": [
            r"(?:gold\s*lock|gold_lock|gold\s*locking|lock\s*value|lock\s*percentage|metal\s*lock|lock)\s*[:\-]?\s*([0-9]+(?:\.[0-9]+)?)",
            r"(?:gold\s*lock|gold_lock|gold\s*locking|lock\s*value|lock\s*percentage|metal\s*lock|lock)\s*[:\-]?\s*[\n\r]+?\s*([0-9]+(?:\.[0-9]+)?)",
            r"(?:gold\s*lock|gold_lock|gold\s*locking|lock\s*value|lock\s*percentage|metal\s*lock|lock)\s{2,}([0-9]+(?:\.[0-9]+)?)",
            r"(?:goldlock|gold_lock|lock)\s*([0-9]+(?:\.[0-9]+)?)"
        ],

    }
    
    EMPTY_PO = {
        "po_number": None,
        "customer_name": None,
        "vendor_number": None,
        "po_date": None,
        "delivery_date": None,
        "cancel_date": None,
        "gold_karat": None,
        "ec_style_number": None,
        "customer_style_number": None,
        "color": None,
        "quantity": None,
        "description": None,
        "gold_lock": None,
    }


    @staticmethod
    def normalize_text(text: str) -> str:
        # First, join multiline customer addresses
        text = re.sub(r'Ship\s+Ostbye\s+To\s*:\s*\n', 'Ship Ostbye To: ', text)
        # preserve newlines, normalize spaces only
        text = re.sub(r"[ \t]+", " ", text)
        return text.strip()
    
    def strip_table_headers(self, text: str) -> str:
        HEADER_PATTERNS = [
            r"item\s+description\s+material\s+quantity\s+expected\s+delivery",
            r"no\s+part\s+no\s+description\s+qty\s+unit\s+price\s+total",
            r"description\s+material\s+quantity\s+expected\s+delivery",
        ]

        for pat in HEADER_PATTERNS:
            text = re.sub(pat, "", text, flags=re.IGNORECASE)

        return text

    # ---------------- PO FIELD EXTRACTION ---------------- #
    @staticmethod
    def extract_po_fields_regex(text: str) -> dict:
        out = SharepointService.EMPTY_PO.copy()
        text = SharepointService.normalize_text(text)

        for field, patterns in SharepointService.PO_REGEX_PATTERNS.items():
            for pat in patterns:
                m = re.search(pat, text, re.IGNORECASE | re.MULTILINE | re.DOTALL)
                if m:
                    out[field] = m.group(1) if m.groups() else m.group(0)
                    break

        return out
    
    MANDATORY_FIELDS = ["po_number", "customer_name", "vendor_number", "po_date", "delivery_date", "quantity"]
    
    ITEM_ONLY_FIELDS = {
        "quantity",
        "gold_karat",
        "description"
    }
    
    ITEM_REPEAT_KV_REGEX = re.compile(
            r"""
            Description\s*:\s*(?P<description>.+?)\s*
            Qty\s*:\s*(?P<quantity>\d+)
            """,
            re.IGNORECASE | re.DOTALL | re.VERBOSE
        )
    
    def extract_karat(self, description: str) -> Optional[str]:
        """Extract gold karat from description."""
        if not description:
            return None
        
        # Look for patterns like 22K, 24K, etc.
        karat_match = re.search(r'(\d{2})K', description, re.IGNORECASE)
        if karat_match:
            return karat_match.group(1)
        
        return None

    
    DB_FIELD_LIMITS = {
        "customer_name": 255,
        "vendor_number": 100,
        "po_number": 100,
        "gold_karat": 100,
        "ec_style_number": 100,
        "customer_style_number": 100,
        "color": 50,
        "quantity": 100,
        "gold_lock": 100,
    }
 
    def trim_to_db_limits(self, data: dict) -> dict:
        trimmed = {}
    
        for field, value in data.items():
            if not value:
                trimmed[field] = value
                continue
    
            if field in self.DB_FIELD_LIMITS:
                max_len = self.DB_FIELD_LIMITS[field]
                trimmed[field] = value[:max_len]
            else:
                trimmed[field] = value
    
        return trimmed

    async def extract_po_fields(self, text: str) -> dict:
        regex_data_response = self.extract_po_fields_regex(text)
        regex_data = self.trim_to_db_limits(regex_data_response)
        logger.info(f"regex data: {regex_data}")

        if all(regex_data.get(f) for f in self.MANDATORY_FIELDS) and len(text.strip()) >= 50:
            return regex_data

        if len(text.strip()) < 50:
            logger.info("Skipping LLM — text too short")
            return self.EMPTY_PO

        llm_data = await self.extract_po_fields_from_llm(text)
        logger.info(f"LLM data: {llm_data}")

        # Merge: LLM fills only the fields regex missed — never blanks a regex value
        final = llm_data
        # final = regex_data.copy()
        # for k, v in llm_data.items():
        #     if v not in (None, "", "null", "N/A"):
        #         if not final.get(k):          # only fill empty slots
        #             final[k] = v

        return final if any(final.values()) else self.EMPTY_PO
    
    
    async def extract_po_fields_from_llm(self, text: str) -> dict:
        field_list = json.dumps(self.PO_FIELD_NAMES, indent=2)
        prompt = f"""
    You are a precision extraction engine for Jewelry Purchase Orders.

    Your job is to extract structured data from raw text parsed from files
    (PDFs, Excel sheets, Word docs, images via OCR, or email bodies).

    ===========================
    GLOBAL RULES (NON-NEGOTIABLE)
    ===========================
    - Extract ONLY values that are EXPLICITLY present with a clear label.
    - NEVER infer, guess, or derive values from context.
    - NEVER merge two fields into one value.
    - NEVER include field label names inside field values.
    - If a value is absent, ambiguous, or only implied return null.
    - Preserve original formatting (casing, spacing, units) in extracted values.
    - Do not normalize, rewrite, or clean values.

    ===========================
    SOURCE-SPECIFIC PARSING RULES
    ===========================

    [EXCEL / CSV TEXT]
    - Text arrives as row-by-row dumps separated by pipe characters.
    - First row is typically the header row — match column names to fields.
    - If multiple data rows exist extract the FIRST valid row only (unless items list).
    - Ignore "Sheet: name" prefix lines.

    [PDF / OCR TEXT]
    - Labels and values may be on separate lines.
    - For key:value pairs the value is text immediately AFTER the colon on the same line,
    or the first non-empty line below the label.

    [IMAGE / OCR TEXT]
    - Apply light OCR error tolerance: "P0 Number" likely means "PO Number".
    - Only extract if you are more than 90 percent certain of the value.

    ===========================
    FIELD-BY-FIELD RULES
    ===========================

    PO NUMBER:
    - Labels: "PO Number", "P.O.", "PO#", "PO #", "Purchase Order Number", "Order No"
    - Value must look like an alphanumeric PO identifier e.g. "PO-2024-001", "12345A".
    - If multiple PO numbers appear return the FIRST one.

    CUSTOMER NAME:
    - Labels: "Customer", "Customer Name", "Bill To", "Sold To", "Buyer"
    - Customer Name can also be an alphanumeric code (e.g., DM5-GER, ABC-123, XYZ LTD).
    - If a field labeled "Customer" exists, ALWAYS extract its full value exactly as written.
    - Do NOT extract address lines as the customer name.

    VENDOR NUMBER:
    - Labels: "Vendor", "Vendor No", "Vendor Number", "Supplier Code", "Vendor ID"
    - Value is typically numeric or alphanumeric e.g. "V-1042", "5583".

    PO DATE:
    - Labels: "PO Date", "Order Date", "Purchase Order Date", "Date"
    - Accept bare "Date" only if near PO Number or document title.
    - Return date exactly as written.

    DELIVERY DATE:
    - Labels: "Delivery Date", "Due Date", "Ship By", "Required By", "Expected Delivery"
    - Return null if only per-item delivery dates exist.

    CANCEL DATE:
    - Labels: "Cancel Date", "Cancellation Date", "Cancel By", "Void After"

    GOLD KARAT:
    - Labels: "Gold Karat", "Karat", "Metal", "Gold Purity", "KT", "K"
    - Valid examples: "14KW", "18K", "22K Yellow", "18 KT", "14K Rose Gold"
    - Label MUST be a standalone field header not embedded in a description.

    EC STYLE NUMBER:
    - Labels: "EC Style", "EC Style No", "EC Style Number", "Style#", "IDD Style #"
    - Alphanumeric style code. Do NOT extract customer-side style numbers.

    CUSTOMER STYLE NUMBER:
    - Labels: "Customer Style", "Cust Style", "Your Style", "Vendor Style #"
    - This is the buyer's own reference number distinct from EC Style.

    COLOR:
    - Labels: "Color", "Colour", "Metal Color"
    - Do NOT extract color embedded in description sentences.

    DESCRIPTION:
    - Labels: "Description", "Item Description", "Particulars", "Goods Description"
    - Extract the FULL value exactly as written.

    QUANTITY:
    - Labels: "Quantity", "Qty", "Pieces", "Pcs", "Ord#"
    - Value must be numeric.
    - Do NOT sum multiple item quantities — return null if only per-item quantities exist.

    GOLD LOCK:
    - Labels: "Gold Lock", "Lock Value", "Lock %", "Metal Lock"
    - Value may be a percentage or numeric.

    ===========================
    MULTI-ROW ITEM DETECTION
    ===========================
    If the text contains a TABLE with multiple product/item rows return:
    "has_line_items": true
    Still extract all header-level fields. Return null for per-item fields
    (Delivery Date, Gold Karat, Description, Quantity, Color) in the main record.

    ===========================
    OUTPUT FORMAT
    ===========================
    Return a single valid JSON object with EXACTLY these keys:
    {field_list}
    Plus one optional key: "has_line_items" (boolean, default false).

    Rules:
    - All string values: preserve original casing and formatting.
    - Missing or ambiguous values: null (not empty string, not "N/A").
    - No explanation, preamble, or markdown — JSON only.

    ===========================
    TEXT TO EXTRACT FROM:
    ===========================
    {text}
    """
        try:
            resp = openai_client.chat.completions.create(
                model="gpt-4o-mini",
                temperature=0,
                messages=[{"role": "user", "content": prompt}]
            )
            raw = resp.choices[0].message.content.strip()
            raw = re.sub(r"^```(?:json)?|```$", "", raw, flags=re.MULTILINE | re.IGNORECASE).strip()
            match = re.search(r"\{.*\}", raw, re.DOTALL)
            if not match:
                return self.EMPTY_PO
            data = json.loads(match.group())
            result = {}
            for k in self.PO_FIELD_NAMES:
                v = data.get(k)
                result[k] = v if v not in ("", None, "null", "N/A") else None
            return result
        except Exception as e:
            logger.error(f"LLM extraction failed: {e}")
            return self.EMPTY_PO

        
    def normalize_po_date_ddmmyyyy(self,date_str: Optional[str]) -> Optional[str]:
        """
        Converts LLM or regex date output to YYYY-MM-DD string.
        Returns None if parsing fails.
        """
        if not date_str:
            return None

        date_str = date_str.strip()

        date_formats = [
            "%Y-%m-%d",    # 2025-07-11
            "%d-%m-%Y",    # 11-07-2025
            "%m-%d-%Y",    # 07-11-2025
            "%m/%d/%Y",    # 07/11/2025
            "%d/%m/%Y",    # 11/07/2025
            "%y-%m-%d",    # 25-07-11
            "%m/%d/%y",    # 07/11/25
            "%d/%m/%y",    # 11/07/25
            "%b/%d/%Y",    # Jul/11/2025
            "%B/%d/%Y",    # July/11/2025
            "%d-%b-%Y",    # 11-Jul-2025
            "%d-%B-%Y",    # 11-July-2025
        ]

        for fmt in date_formats:
            try:
                dt = datetime.strptime(date_str, fmt)
                return dt.strftime("%Y-%m-%d")
            except ValueError:
                continue

        return None


    def normalize_attachment_text(self, text: str) -> str:
        if not text:
            return ""
        text = text.replace("\u2013", "-").replace("\u2014", "-")
        text = re.sub(r"\s*-\s*", "-", text)
        # DO NOT flatten newlines — they preserve table row/column structure
        # Fix broken word-wrap across lines (word- \n word → word - word)
        text = re.sub(r"([A-Za-z])-\n\s*([A-Za-z])", r"\1-\2", text)
        # Normalize dates like 2025-07-06
        text = re.sub(
            r"(\d{4})-(\d{1,2})-(\d{1,2})",
            lambda m: f"{m.group(1)}-{int(m.group(2)):02d}-{int(m.group(3)):02d}",
            text,
        )
        text = re.sub(r"[ \t]+", " ", text)       # collapse horizontal whitespace only
        text = re.sub(r"\n{3,}", "\n\n", text)    # max 2 consecutive blank lines
        return text.strip()

    
    ITEM_ROW_REGEX = re.compile(
            r"""
            ^\s*
            (?P<row_no>\d+)\s+
            (?P<part_no>[A-Z0-9\-]+)\s+
            (?P<description>.+?)\s+
            (?P<quantity>\d+)\s+
            (?P<unit_price>\d+(?:\.\d+)?)\s+
            (?P<total>\d+(?:\.\d+)?)
            \s*$
            """,
            re.IGNORECASE | re.MULTILINE | re.VERBOSE
        )
    
    
    PO_REPEAT_BLOCK_REGEX = re.compile(
            r"""
            P\.?O\.?\s*Number\s*:\s*(?P<po_number>[A-Z0-9\-]+)\s*
            P\.?O\.?\s*Date\s*:\s*(?P<po_date>[0-9]{2}-[A-Za-z]{3}-[0-9]{4})\s*
            Vendor\s*:\s*(?P<vendor>.+?)\s*
            Description\s*:\s*(?P<description>.+?)\s*
            Qty\s*:\s*(?P<quantity>\d+)
            """,
            re.IGNORECASE | re.VERBOSE | re.DOTALL
        )
    
    ITEM_BLOCK_REGEX = re.compile(
            r"""
            Expected\s*Delivery\s*:\s*(?P<delivery_date>\d{4}-\d{2}-\d{2})
            .*?
            Item\s*Description\s*:\s*(?P<description>.+?)
            .*?
            Material\s*:\s*(?P<material>.+?)
            .*?
            Quantity\s*:\s*(?P<quantity>\d+)
            """,
            re.IGNORECASE | re.DOTALL | re.VERBOSE
        )
    
    ITEM_TABLE_REGEX = re.compile(
            r"""
            (?P<sku>[A-Z0-9]{3,})\s+
            (?P<description>[A-Za-z\s\-]+?)\s+
            (?P<gold_karat>\d{2})K\s+
            (?P<quantity>\d+)
            """,
            re.IGNORECASE | re.VERBOSE
        )
    
    ITEM_COLUMN_REGEX = re.compile(
            r"""
            (?P<description>(?!expected\s+delivery)[A-Za-z][A-Za-z\s\-–]+?)\s+
            (?P<material>\d{2}K\s*Gold(?:\s*\+\s*Diamond)?)\s+
            (?P<quantity>\d+)\s+
            (?P<delivery_date>\d{4}-\d{2}-\d{2})
            """,
            re.IGNORECASE | re.VERBOSE
        )
    
    ITEM_PIPE_TABLE_REGEX = re.compile(
            r"""
            (?P<description>[^|]+)\s*\|\s*
            (?P<material>\d{2}K\s*Gold(?:\s*\+\s*Diamond)?)\s*\|\s*
            (?P<quantity>\d+)\s*\|\s*
            (?P<delivery_date>\d{4}-\d{2}-\d{2})
            """,
            re.IGNORECASE | re.VERBOSE
        )
    
    PO_BLOCK_REGEX = re.compile(
            r"""
            P\.?O\.?\s*Number\s*:\s*(?P<po_number>[^\n]+)
            .*?
            P\.?O\.?\s*Date\s*:\s*(?P<po_date>[^\n]+)
            .*?
            Vendor\s*:\s*(?P<vendor>[^\n]+)
            .*?
            Description\s*:\s*(?P<description>[^\n]+)
            .*?
            Qty\s*:\s*(?P<quantity>\d+)
            """,
            re.IGNORECASE | re.DOTALL | re.VERBOSE
        )

    
    def strip_item_sections(self, text: str) -> str:
        return re.sub(
            r"Item\s+Description\s+Material\s+Quantity\s+Expected\s+Delivery",
            "",
            text,
            flags=re.IGNORECASE
        )
        
        
    def extract_purchase_order_table(self, text: str):
        """Extract data from purchase order table format like in your example."""
        results = []
        
        # Extract common header data
        po_number_match = re.search(r'P\.O\.\s*Number\s*:\s*(JG-PO-2025-0043)', text, re.IGNORECASE)
        po_date_match = re.search(r'P\.O\.\s*Date\s*:\s*(11-Jul-2025)', text, re.IGNORECASE)
        vendor_match = re.search(r'Vendor\s*:\s*(.+?)(?=\s+Ship\s+To:)', text, re.IGNORECASE)
        
        po_number = po_number_match.group(1).strip() if po_number_match else None
        po_date = po_date_match.group(1).strip() if po_date_match else None
        vendor = vendor_match.group(1).strip() if vendor_match else None
        
        # Look for the table in the text
        lines = text.split('\n')
        in_table = False
        
        for line in lines:
            line = line.strip()
            
            # Look for table header
            if 'NoPart No Description Qty Unit Price Total' in line:
                in_table = True
                continue
            
            if in_table:
                # Stop when we hit totals
                if 'Sub-Total' in line or 'Total' in line or 'Tax' in line:
                    break
                
                # Try to parse table row
                # Your text shows: "1JWL001 22K Gold Necklace - Antique Finish 10 75000.00 750000.00"
                
                # Pattern 1: Look for JWL codes
                if 'JWL' in line:
                    # Split by spaces and try to find quantity
                    parts = line.split()
                    
                    # Look for quantity - it's a number followed by price pattern
                    for i, part in enumerate(parts):
                        if part.isdigit() and i + 1 < len(parts) and '.' in parts[i + 1]:
                            quantity = int(part)
                            # Reconstruct description
                            description_parts = []
                            
                            # Start from position 1 (skip row number + JWL code)
                            # Find where the JWL code is
                            for j in range(len(parts)):
                                if 'JWL' in parts[j]:
                                    # Start description after JWL code
                                    k = j + 1
                                    while k < i:  # Until we reach quantity
                                        description_parts.append(parts[k])
                                        k += 1
                                    break
                            
                            description = ' '.join(description_parts)
                            
                            if description:
                                results.append({
                                    "po_number": po_number,
                                    "po_date": po_date,
                                    "vendor": vendor,
                                    "description": description,
                                    "quantity": quantity
                                })
                            break
        
        return results
    
    
    ITEM_INLINE_REGEX = re.compile(
        r"""
        (?P<description>[A-Za-z0-9\s\-:.]+?)
        \s+
        (?P<quantity>\d+(?:\.\d+)?)
        \s+EA
        """,
        re.IGNORECASE | re.VERBOSE
    )
    
        
    @staticmethod  
    def clean_item_description(desc: str) -> str:
        desc = re.sub(r"^[A-Z0-9.\-/]+\s+", "", desc)
        desc = re.sub(r"\b\d{2}K[W]?\b", "", desc)
        return re.sub(r"\s+", " ", desc).strip()



    async def extract_po_header(self,text: str):
        return await self.extract_po_fields(text)
    
    IMAGE_EXTENSIONS = (".png", ".jpg", ".jpeg", ".bmp", ".tiff", ".webp")
    
    def extract_relative_folder_path(self, graph_path: str) -> str:
        """
        Converts:
        /drives/{id}/root:/Elegant Collection Software/SubFolder
        → Elegant Collection Software/SubFolder
        """
        if not graph_path:
            return ""

        # Everything after 'root:/'
        if "root:/" in graph_path:
            return graph_path.split("root:/", 1)[1]

        return graph_path


    # ---------------- MAIN FLOW ---------------- #
    async def fetch_and_save_sharepoint_files(
        self,
        access_token: str,
        user_id: int,
        folders: list[str],
        from_date: str,
        to_date: str,
        site_id: str,
    ):
        saved, failed = [], []
        extracted_sharepoint_po_ids: list[int] = []

        drive_id = await self.get_drive_id(access_token, site_id)
        keywords = await self.sp_repo.fetch_keywords()

        async with aiohttp.ClientSession(
            headers={"Authorization": f"Bearer {access_token}"}
        ) as session:

            for folder in folders or [""]:
                files = await self.fetch_drive_files(
                    access_token, drive_id, folder, from_date, to_date
                )

                for f in files:
                    try:
                        url = f.get("@microsoft.graph.downloadUrl")
                        if not url:
                            continue

                        async with session.get(url) as r:
                            data = await r.read()

                        # ---------------- HASH CHECK ----------------
                        file_hash = await self.generate_file_hash(data)
                        if await self.sp_repo.file_exists(user_id, file_hash):
                            continue

                        # ---------------- TEXT EXTRACTION ----------------
                        text = await self.extract_text_from_bytes(
                            data, f["name"], f["file"]["mimeType"]
                        )

                        if not text:
                            continue

                        # ---------------- Normalize ONCE (CRITICAL) ----------------
                        normalized_text = self.normalize_attachment_text(text)

                        # ---------------- KEYWORD DETECTION ----------------
                        matched_keywords, _ = await self.detect_keywords(
                            normalized_text, keywords
                        )
                        if not matched_keywords:
                            continue

                        # ---------------- SAVE FILE ----------------
                        parent_path = f.get("parentReference", {}).get("path", "")
                        folder_name = self.extract_relative_folder_path(parent_path)

                        sharepoint_file_id = await self.sp_repo.save_sharepoint_file(
                            user_id=user_id,
                            file_name=f["name"],
                            file_type=f["file"]["mimeType"],
                            file_path=f["webUrl"],
                            file_size=f.get("size", 0),
                            folder_name=folder_name,
                            uploaded_on=self.graph_datetime_to_mysql(
                                f["createdDateTime"]
                            ),
                            file_hash=file_hash,
                            created_by=user_id,
                        )

                        # ---------------- HEADER EXTRACTION ----------------
                        header = await self.extract_po_fields(normalized_text)

                        items = None #for now we are showing only header level data, item level data extraction will be added in next phase

                        # ---------------- HEADER ONLY ----------------
                        if not items:
                            if header.get("po_number") or header.get("customer_name"):
                                po_det_id = await self.sp_repo.insert_sharepoint_po_details(
                                    sharepoint_file_id=sharepoint_file_id,
                                    user_id=user_id,
                                    po_number=header.get("po_number"),
                                    customer_name=header.get("customer_name"),
                                    vendor_number=header.get("vendor_number"),
                                    po_date=self.normalize_po_date_ddmmyyyy(header.get("po_date")),
                                    delivery_date=self.normalize_po_date_ddmmyyyy(header.get("delivery_date")),
                                    cancel_date=self.normalize_po_date_ddmmyyyy(header.get("cancel_date")),
                                    gold_karat=header.get("gold_karat"),
                                    ec_style_number=header.get("ec_style_number"),
                                    customer_style_number=header.get("customer_style_number"),
                                    color=header.get("color"),
                                    quantity=header.get("quantity"),
                                    description=header.get("description"),
                                    created_by=user_id,
                                    gold_lock=header.get("gold_lock"),
                                )
                                extracted_sharepoint_po_ids.append(po_det_id)

                        # ---------------- MULTI ITEMS ----------------
                        else:
                            for item in items:
                                po_number = item.get("po_number") or header.get("po_number")
                                po_date = item.get("po_date") or header.get("po_date")
                                vendor = item.get("vendor") or header.get("vendor_number")

                                po_det_id = await self.sp_repo.insert_sharepoint_po_details(
                                    sharepoint_file_id=sharepoint_file_id,
                                    user_id=user_id,
                                    po_number=po_number,
                                    customer_name=header.get("customer_name"),
                                    vendor_number=vendor,
                                    po_date=self.normalize_po_date_ddmmyyyy(po_date),
                                    delivery_date=self.normalize_po_date_ddmmyyyy(
                                        item.get("delivery_date")
                                    ) or self.normalize_po_date_ddmmyyyy(
                                        header.get("delivery_date")
                                    ),
                                    cancel_date=self.normalize_po_date_ddmmyyyy(
                                        header.get("cancel_date")
                                    ),
                                    gold_karat=item.get("gold_karat"),
                                    ec_style_number=header.get("ec_style_number"),
                                    customer_style_number=header.get("customer_style_number"),
                                    color=header.get("color"),
                                    quantity=item.get("quantity"),
                                    description=item.get("description"),
                                    created_by=user_id,
                                    gold_lock=item.get("gold_lock"),
                                )
                                extracted_sharepoint_po_ids.append(po_det_id)

                        saved.append(f["name"])

                    except Exception as e:
                        failed.append({
                            "file": f.get("name"),
                            "error": str(e),
                        })

        return {
            "saved_count": len(saved),
            "failed_count": len(failed),
            "saved_files": saved,
            "failed_files": failed,
            "extracted_sharepoint_po_ids": extracted_sharepoint_po_ids,
        }
    #--------------------------data comparison logic start--------------------------#

    # -------------------------- FIELDS TO COMPARE -------------------------- #
    FIELDS_TO_COMPARE = [
        "customer_name",
        "vendor_number",
        "po_date",
        "po_number",
        "delivery_date",
        "cancel_date",
        "gold_lock",
        "ec_style_number",
        "customer_style_number",
        "gold_karat",
        "color",
        "quantity",
        "description",
    ]

    # -------------------------- LLM FALLBACK -------------------------- #
    @staticmethod
    async def llm_batch_match(scanned_pos, system_pos):
        prompt = f"""
    You are a PO matching engine.

    Match scanned POs to system POs using ONLY:
    - customer_name
    - po_number

    Rules:
    - Handle spelling mistakes, abbreviations, extra/missing letters.
    - One scanned PO matches at most one system PO.
    - If no confident match exists, return null.

    Return ONLY JSON:
    [
    {{
        "scanned_po_det_id": number,
        "system_po_id": number | null,
        "confidence": 0.0-1.0
    }}
    ]

    Scanned POs:
    {json.dumps(scanned_pos)}

    System POs:
    {json.dumps(system_pos)}
    """

        resp = openai_client.chat.completions.create(
            model="gpt-4.1-mini",
            temperature=0,
            messages=[{"role": "user", "content": prompt}]
        )

        raw = resp.choices[0].message.content.strip()
        raw = re.sub(r"^```json|```$", "", raw, flags=re.IGNORECASE).strip()
        return json.loads(raw)


    async def llm_batch_compare(self,matched_pairs):
        prompt = f"""
    You are an expert PO field comparison engine.

    Compare ONLY the following fields:
    {self.FIELDS_TO_COMPARE}

    Your goal is to detect **real business mismatches**. Do NOT report differences caused by minor spelling mistakes, abbreviations, word order, or formatting.

    Rules:
    1. Treat minor spelling errors, missing/extra letters, or phonetic variations as SAME.
    2. Treat abbreviations and expansions (e.g., Pvt, Ltd, Private Limited) as SAME.
    3. Ignore punctuation, dots, commas, extra spaces, capitalization.
    4. Ignore word order in names, colors, or descriptions.
    5. Ignore formatting differences in dates (YYYY-MM-DD, DD/MM/YYYY) or numbers/quantities.
    6. Only report a mismatch if the values clearly indicate different meanings or entities.

    Return ONLY JSON in the following format:
    [
    {{
        "sharepoint_po_det_id": number,
        "system_po_id": number,
        "field": string,
        "scanned_value": string,
        "system_value": string
    }}
    ]

    Here are the matched PO pairs to compare:
    {json.dumps(matched_pairs, indent=2)}
    """

        # Call OpenAI LLM
        resp = openai_client.chat.completions.create(
            model="gpt-4.1-mini",
            temperature=0,
            messages=[{"role": "user", "content": prompt}]
        )

        # Clean up response
        raw = resp.choices[0].message.content.strip()
        raw = re.sub(r"^```json|```$", "", raw, flags=re.IGNORECASE).strip()

        # Parse JSON
        try:
            result = json.loads(raw)
            if not isinstance(result, list):
                raise ValueError("LLM response is not a JSON list")
            return result
        except Exception as e:
            # fallback: return empty list if parsing fails
            print(f"LLM parsing error: {e}")
            return []

    @staticmethod
    def chunk(data, size):
        for i in range(0, len(data), size):
            yield data[i:i + size]

    
    #---------------------------Table Data ----------------------------
    async def missing_po_data_fetch(request: Request, frontendRequest):
        data = await SharepointRepo.fetch_missing_po_data(request, frontendRequest)
        # FIX: Return empty list if None, and return the LIST directly (no wrapper object)
        return data if data else []
            
    async def mismatch_po_data_fetch(request: Request, frontendRequest):
        data = await SharepointRepo.fetch_mismatch_po_data(request, frontendRequest)
        return data if data else []
            
    async def matched_po_data_fetch(request: Request, frontendRequest):
        data = await SharepointRepo.fetch_matched_po_data(request, frontendRequest)
        return data if data else []
    
    
    #Download Missing Report and Mismatch Report
    async def download_sharepoint_missing_po_report(
        request: Request,
        user_id: int,
        role_id: int,
        format: str,
        selected_ids: Optional[List[int]] = None
    ):
        selected_ids = selected_ids or []
        data = await SharepointRepo.download_sharepoint_missing_po_report(request, user_id, role_id,selected_ids)

        if not data:
            raise HTTPException(status_code=404, detail="No missing PO data available")

        df = pd.DataFrame(data)
        filename_prefix = "po_missing_report"

        if format == "excel":
            output = io.BytesIO()
            df.to_excel(output, index=False)
            output.seek(0)

            return (
                output,
                f"{filename_prefix}.xlsx",
                "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
            )

        elif format == "pdf":
            return SharepointService._generate_po_pdf(df, filename_prefix)

        else:
            raise HTTPException(status_code=400, detail="Invalid file format")

    async def download_sharepoint_mismatch_po_report(
        request: Request,
        user_id: int,
        role_id: int,
        format: str,
        selected_ids: Optional[List[int]] = None
    ):
        selected_ids = selected_ids or []
        data = await SharepointRepo.download_sharepoint_mismatch_po_report(request, user_id, role_id,selected_ids)

        if not data:
            raise HTTPException(status_code=404, detail="No mismatch PO data available")

        df = pd.DataFrame(data)
        filename_prefix = "po_mismatch_report"

        if format == "excel":
            output = io.BytesIO()
            df.to_excel(output, index=False)
            output.seek(0)

            return (
                output,
                f"{filename_prefix}.xlsx",
                "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
            )

        elif format == "pdf":
            return SharepointService._generate_po_pdf(df, filename_prefix)

        else:
            raise HTTPException(status_code=400, detail="Invalid file format")
        
        
    async def download_selected_po_report(
        request: Request,
        user_id: int,
        role_id: int,
        sharepoint_missing_ids: List[int] = None,
        sharepoint_mismatch_ids: List[int] = None,
        sharepoint_matched_ids: List[int] = None,
        format: str = "excel"
    ):
        data = await SharepointRepo.download_selected_po_report(
            request, user_id, role_id, sharepoint_missing_ids, sharepoint_mismatch_ids, sharepoint_matched_ids
        )

        if not data:
            return None, "no_data.xlsx" if format=="excel" else "no_data.pdf", "application/octet-stream"

        df = pd.DataFrame(data)
        filename_prefix = "sharepoint_po_report"

        if format == "excel":
            output = io.BytesIO()
            df.to_excel(output, index=False)
            output.seek(0)
            return (
                output, 
                f"{filename_prefix}.xlsx", 
                "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
            )
            
        elif format == "pdf":
            return SharepointService._generate_po_pdf(df, filename_prefix)

        else:
            raise HTTPException(status_code=400, detail="Invalid file format")

        
    # #Last Sync On Dashboard(Sharepoint)
    async def get_last_sync_by_user_id(user_id: int,role_id: int,request: Request):
        try:
            last_sync_data = await SharepointRepo.get_last_sync_by_user_id(user_id,role_id,request)
            return last_sync_data
        except Exception as e:
            raise Exception(f"Error fetching last sync data: {str(e)}")   
        
    #Adding and Update comment for po missing and po mismatch from UI
    async def save_sharepoint_po_comment(
        report_type: str,
        record_id: int,
        comment: str,
        request: Request
    ):
        if report_type == "missing":
            return await SharepointRepo.save_sharepoint_po_missing_comment(
                record_id, comment, request
            )

        elif report_type == "mismatch":
            return await SharepointRepo.save_sharepoint_po_mismatch_comment(
                record_id, comment, request
            )


    #For Fetching the PO comment ON UI 
    async def fetch_sharepoint_po_comment(
            report_type: str,
            record_id: int,
            request: Request
        ) -> str | None:

            if report_type == "missing":
                return await SharepointRepo.fetch_sharepoint_missing_po_comment(
                    record_id, request
                )

            elif report_type == "mismatch":
                return await SharepointRepo.fetch_sharepoint_mismatch_po_comment(
                    record_id, request
                )

            else:
                raise ValueError("Invalid report type")


    #For Ignoring the PO in Next Sync On UI
    async def ignore_sharepoint_po(
            report_type: str,
            record_id: int,
            request: Request
        ) -> bool:

            if report_type == "missing":
                return await SharepointRepo.ignore_sharepoint_missing_po(
                    record_id, request
                )

            elif report_type == "mismatch":
                return await SharepointRepo.ignore_sharepoint_mismatch_po(
                    record_id, request
                )

            else:
                raise ValueError("Invalid report type")
            

    # ============================================================
    # data comparison logic start
    # ============================================================
    FIELDS_TO_COMPARE = [
        "customer_name",
        "vendor_number",
        "po_date",
        "po_number",
        "delivery_date",
        "cancel_date",
        "gold_lock",
        "ec_style_number",
        "customer_style_number",
        "gold_karat",
        "color",
        "quantity",
        "description",
    ]

    @staticmethod
    def make_json_safe(obj):
        if isinstance(obj, (date, datetime)):
            return obj.isoformat()
        if isinstance(obj, Decimal):
            return float(obj)
        if isinstance(obj, bytes):
            try:
                return obj.decode("utf-8")
            except UnicodeDecodeError:
                return base64.b64encode(obj).decode("utf-8")
        return obj


    async def llm_batch_match(scanned_pos, system_pos):
        prompt = f"""
    You are a PO matching engine.

    Match scanned POs to system POs using ONLY:
    - customer_name
    - po_number

    Rules:
    - Handle spelling mistakes, abbreviations, extra/missing letters.
    - One scanned PO matches at most one system PO.
    - If no confident match exists, return null.

    Return ONLY JSON:
    [
    {{
        "scanned_po_det_id": number,
        "system_po_id": number | null,
        "confidence": 0.0-1.0
    }}
    ]

    Scanned POs:
    {json.dumps(scanned_pos)}

    System POs:
    {json.dumps(system_pos)}
    """

        resp = openai_client.chat.completions.create(
            model="gpt-4.1-mini",
            temperature=0,
            messages=[{"role": "user", "content": prompt}]
        )

        raw = resp.choices[0].message.content.strip()
        raw = re.sub(r"^```json|```$", "", raw, flags=re.IGNORECASE).strip()
        return json.loads(raw)


    async def llm_batch_compare(self, pairs_for_llm):

        prompt = f"""
    You are an expert PO field comparison engine.

    Determine if the scanned value and system value represent the SAME meaning.

    Rules:
    - Treat spelling mistakes as SAME
    - Treat abbreviations as SAME
    - Treat casing differences as SAME
    - Treat punctuation differences as SAME
    - Treat word order differences as SAME

    Examples considered SAME:
    "yellow color" vs "color yellow"
    "intl corp" vs "international corporation"
    "ring size 7" vs "size 7 ring"

    Only report mismatches if the meaning is clearly different.

    Return ONLY JSON:
    [
    {{
        "sharepoint_po_det_id": number,
        "system_po_id": number,
        "field": string,
        "scanned_value": string,
        "system_value": string
    }}
    ]

    PO field data:
    {json.dumps(pairs_for_llm, indent=2)}
    """

        resp = openai_client.chat.completions.create(
            model="gpt-4.1-mini",
            temperature=0,
            messages=[{"role": "user", "content": prompt}]
        )

        raw = resp.choices[0].message.content.strip()
        raw = re.sub(r"^```json|```$", "", raw, flags=re.IGNORECASE).strip()

        try:
            result = json.loads(raw)
            return result if isinstance(result, list) else []
        except Exception as e:
            print(f"LLM parsing error: {e}")
            return []
        

    @staticmethod
    def chunk(data, size):
        for i in range(0, len(data), size):
            yield data[i:i + size]

    @staticmethod
    def normalize_po(po):
        if not po:
            return ""
        return re.sub(r'[^A-Za-z0-9]', '', po).lower()


    def normalize_value(self, v):
        if not v:
            return ""

        v = str(v).lower().strip()

        # remove company suffix
        v = re.sub(r'\b(pvt|ltd|limited|private|corp|inc)\b', '', v)

        # remove special chars
        v = re.sub(r'[^a-z0-9\s]', '', v)

        # normalize spaces
        v = re.sub(r'\s+', ' ', v)

        return v.strip()


    def is_field_match(self, field, s_val, t_val):
        if s_val in (None, "") or t_val in (None, ""):
            return False

        s = self.normalize_value(s_val)
        t = self.normalize_value(t_val)

        # EXACT
        if s == t:
            return True

        # TEXT FIELDS
        if field in ["customer_name", "description", "color"]:
            return fuzz.token_sort_ratio(s, t) >= 85

        # DATE FIELDS
        if field in ["po_date", "delivery_date"]:
            try:
                d1 = pd.to_datetime(s_val, errors="coerce")
                d2 = pd.to_datetime(t_val, errors="coerce")

                if d1 is not None and d2 is not None:
                    return abs((d1 - d2).days) <= 2
            except:
                pass

        # NUMERIC
        if field in ["quantity", "gold_lock"]:
            try:
                return float(s_val) == float(t_val)
            except:
                return False

        return False
    

    def get_field_mismatches(self, scanned, system):
        mismatches = {}

        for field in self.FIELDS_TO_COMPARE:

            if not self.is_field_match(field, scanned.get(field), system.get(field)):
                mismatches[field] = {
                    "scanned": scanned.get(field),
                    "system": system.get(field)
                }

        return mismatches


    def find_best_system_match(self, scanned, candidates):

        best_candidate = None
        best_score = -1

        for system in candidates:
            score = 0

            for field in FIELDS_TO_COMPARE:

                scanned_val = self.normalize_value(scanned.get(field))
                system_val = self.normalize_value(system.get(field))

                if not scanned_val or not system_val:
                    continue

                if scanned_val == system_val:
                    score += 1

            if score > best_score:
                best_score = score
                best_candidate = system

        return best_candidate

    # ============================================================
    # PO Recomparison start
    # ============================================================

    # ----------------STABLE SYSTEM PO ID----------------------------
    def make_stable_system_po_id(self, po: dict) -> int:
        raw = "|".join([
            self.normalize_po(str(po.get("po_number") or "")),
            self.normalize_value(str(po.get("customer_number") or "")),
        ])
        full_hash = hashlib.sha256(raw.encode()).digest()
        return int.from_bytes(full_hash[:4], byteorder="big")


    # -------------------------- PO Recomparison ENTRY POINT ------------------------------
    def is_strong_match(self,scanned, system):
        return (
            self.normalize_po(scanned.get("po_number")) == self.normalize_po(system.get("po_number"))
            and self.normalize_value(scanned.get("customer_number")) == self.normalize_value(system.get("customer_number"))
        )


    async def reconcile_all_sharepoint_pos(self,user_id, system_pos, sp_repo=SharepointRepo, batch_size=500):
        if sp_repo is None:
            logger.error("[RECONCILE] sp_repo not provided to reconcile_all_sharepoint_pos")
            return {"errors": ["sp_repo missing"]}

        stats = {
            "mismatch_checked": 0,
            "mismatch_resolved": 0,
            "missing_checked": 0,
            "missing_resolved": 0,
            "errors": []
        }

        try:
            # Build system lookup
            system_po_map = defaultdict(list)
            for po in system_pos:
                key = self.normalize_po(po.get("po_number"))
                if key:
                    system_po_map[key].append(po)

            # Run mismatch reconcile — isolated individually 
            try:
                mismatch_stats = await self.reconcile_mismatches(user_id, system_po_map, batch_size, sp_repo)
                stats.update(mismatch_stats)
            except Exception as e:
                logger.error(f"[RECONCILE] reconcile_mismatches failed: {e}", exc_info=True)
                stats["errors"].append(f"mismatch: {str(e)}")

            # Run missing reconcile — isolated individually 
            try:
                missing_stats = await self.reconcile_missing(user_id, system_po_map, batch_size, sp_repo)
                stats.update(missing_stats)
            except Exception as e:
                logger.error(f"[RECONCILE] reconcile_missing failed: {e}", exc_info=True)
                stats["errors"].append(f"missing: {str(e)}")

        except Exception as e:
            logger.error(f"[RECONCILE] reconcile_all_pos outer exception: {e}", exc_info=True)
            stats["errors"].append(str(e))

        return stats


    # ============================================================
    # MISMATCH RECONCILIATION
    # ============================================================
    async def reconcile_mismatches(self, user_id, system_po_map, batch_size, sp_repo=SharepointRepo):

        stats = {"mismatch_checked": 0, "mismatch_resolved": 0}

        active = await sp_repo.get_all_active_sharepoint_mismatches(user_id)
        if not active:
            return stats

        stats["mismatch_checked"] = len(active)

        # Group mismatches by PO
        grouped = defaultdict(list)
        for row in active:
            grouped[row["sharepoint_po_det_id"]].append(row)

        sharepoint_po_det_ids = list(grouped.keys())

        scanned_list = await sp_repo.get_sharepoint_po_details_by_ids(sharepoint_po_det_ids)
        scanned_map = {p["sharepoint_po_det_id"]: p for p in scanned_list}

        for batch_keys in self.chunk(sharepoint_po_det_ids, batch_size):

            for sharepoint_po_det_id in batch_keys:

                scanned = scanned_map.get(sharepoint_po_det_id)
                if not scanned:
                    continue

                candidates = system_po_map.get(self.normalize_po(scanned.get("po_number")))
                if not candidates:
                    continue

                # Find best system match
                system = self.find_best_system_match(scanned, candidates)
                if not system:
                    continue

                rows = grouped[sharepoint_po_det_id]

                # FIELD RESOLUTION CHECK (SMART MATCH)
                all_resolved = True

                for row in rows:
                    field = row["mismatch_attribute"]

                    s_val = scanned.get(field)
                    t_val = system.get(field)

                    if not s_val or not t_val:
                        continue

                    if self.normalize_value(s_val) != self.normalize_value(t_val):
                        all_resolved = False
                        break

                # if still mismatch → skip
                if not all_resolved or not self.is_strong_match(scanned, system):
                    continue


                try:
                    # FIXED COLUMN NAME
                    mismatch_ids = [r["sharepoint_po_mismatch_id"] for r in rows]

                    await sp_repo.deactivate_sharepoint_mismatches(user_id, mismatch_ids)

                    exists = await sp_repo.matched_sharepoint_po_exists(
                        user_id, sharepoint_po_det_id, system["system_po_id"]
                    )

                    if not exists:
                        await sp_repo.insert_sharepoint_matched_po(
                            sharepoint_po_det_id=sharepoint_po_det_id,
                            system_po_id=system["system_po_id"],
                            sharepoint_file_id=scanned.get("sharepoint_file_id"),
                            user_id=user_id,
                            po_number=scanned.get("po_number"),
                            po_date=scanned.get("po_date"),
                            vendor_number=scanned.get("vendor_number"),
                            customer_name=scanned.get("customer_name"),
                            created_by="reconciliation"
                        )

                    stats["mismatch_resolved"] += 1

                except Exception as e:
                    print(f"Mismatch resolve error: {e}")

        return stats


    # ============================================================
    # MISSING RECONCILIATION
    # ============================================================
    def get_field_mismatches(self, scanned, system):
        mismatches = {}

        if scanned.get("po_date") != system.get("po_date"):
            mismatches["po_date"] = {
                "scanned": scanned.get("po_date"),
                "system": system.get("po_date")
            }

        if scanned.get("vendor_number") != system.get("vendor_number"):
            mismatches["vendor_number"] = {
                "scanned": scanned.get("vendor_number"),
                "system": system.get("vendor_number")
            }

        if scanned.get("customer_name") != system.get("customer_name"):
            mismatches["customer_name"] = {
                "scanned": scanned.get("customer_name"),
                "system": system.get("customer_name")
            }

        return mismatches


    async def reconcile_missing(self, user_id, system_po_map, batch_size, sp_repo=SharepointRepo):
        stats = {
            "missing_checked": 0,
            "missing_resolved": 0,
            "mismatch_created": 0
        }

        active = await sp_repo.get_all_active_sharepoint_missing(user_id)
        if not active:
            return stats

        stats["missing_checked"] = len(active)

        ids = list({r["sharepoint_po_det_id"] for r in active if r.get("sharepoint_po_det_id")})

        scanned_list = await sp_repo.get_sharepoint_po_details_by_ids(ids)
        scanned_map = {p["sharepoint_po_det_id"]: p for p in scanned_list}

        for batch in self.chunk(active, batch_size):

            for row in batch:

                po_det_id = row["sharepoint_po_det_id"]
                missing_id = row["sharepoint_po_missing_id"]

                scanned = scanned_map.get(po_det_id)
                if not scanned:
                    continue

                po_number = scanned.get("po_number") or row.get("scanned_value")
                if not po_number:
                    continue

                candidates = system_po_map.get(self.normalize_po(po_number))
                if not candidates:
                    continue

                system = self.find_best_system_match(scanned, candidates)
                if not system:
                    continue

                # ============================================================
                # BASE MATCH (PO + CUSTOMER )
                # ============================================================
                base_match = (
                    self.normalize_po(scanned.get("po_number")) == self.normalize_po(system.get("po_number"))
                    and self.normalize_value(scanned.get("customer_name")) == self.normalize_value(system.get("customer_name"))
                )

                if not base_match:
                    continue

                try:
                    # ============================================================
                    # CHECK ALL FIELDS
                    # ============================================================
                    mismatches = []

                    for field in self.FIELDS_TO_COMPARE:

                        s_val = scanned.get(field)
                        t_val = system.get(field)

                        s_norm = self.normalize_value(str(s_val or ""))
                        t_norm = self.normalize_value(str(t_val or ""))

                        if s_norm != t_norm:
                            mismatches.append({
                                "field": field,
                                "scanned": s_val,
                                "system": t_val
                            })

                    # ============================================================
                    # FULL MATCH
                    # ============================================================
                    if not mismatches:

                        await sp_repo.deactivate_sharepoint_missing_pos(user_id, missing_id)

                        exists = await sp_repo.matched_sharepoint_po_exists(
                            user_id, po_det_id, system["system_po_id"]
                        )

                        if not exists:
                            await sp_repo.insert_sharepoint_matched_po(
                                sharepoint_po_det_id=po_det_id,
                                system_po_id=system["system_po_id"],
                                sharepoint_file_id=scanned.get("sharepoint_file_id"),
                                user_id=user_id,
                                po_number=scanned.get("po_number"),
                                po_date=scanned.get("po_date"),
                                vendor_number=scanned.get("vendor_number"),
                                customer_name=scanned.get("customer_name"),
                                created_by="reconciliation"
                            )

                        stats["missing_resolved"] += 1
                        continue

                    # ============================================================
                    # PARTIAL MATCH → INSERT MISMATCH (ALL 13 FIELDS)
                    # ============================================================

                    mismatch_inserted = False

                    for mm in mismatches:

                        exists = await sp_repo.mismatch_exists(
                            user_id=user_id,
                            sharepoint_po_det_id=po_det_id,
                            system_po_id=system["system_po_id"],
                            mismatch_attribute=mm["field"],
                            scanned_value=str(mm["scanned"] or ""),
                            system_value=str(mm["system"] or "")
                        )

                        if not exists:
                            await sp_repo.insert_mismatch(
                                sharepoint_po_det_id=po_det_id,
                                user_id=user_id,
                                system_po_id=system["system_po_id"],
                                field=mm["field"],
                                scanned_value=str(mm["scanned"] or ""),
                                system_value=str(mm["system"] or ""),
                                comment=f"{mm['field']} mismatch"
                            )

                            mismatch_inserted = True

                    # deactivate only if mismatch inserted
                    if mismatch_inserted:
                        await sp_repo.deactivate_sharepoint_missing_pos(user_id, missing_id)
                        stats["mismatch_created"] += 1

                except Exception as e:
                    print(f"Reconcile error: {e}")

        return stats
    # ============================================================
    # PO Recomparison end
    # ============================================================

    async def fetch_Sp_system_pos_with_oldest_date(self, sp_repo, app):
        oldest_date = await sp_repo.get_sharepoint_oldest_report_date()

        if oldest_date:
            system_pos = await MSSQLRepo.get_po_list(app, oldest_date)
        else:
            system_pos = await MSSQLRepo.get_po_list_without_oldest_date(app)

        # ---- imaginary PK for system_pos ---- #
        for po in system_pos:
            po["system_po_id"] = self.make_stable_system_po_id(po)

        system_pos = [
            {k: self.make_json_safe(v) for k, v in po.items()}
            for po in system_pos
        ]

        return system_pos

    # ============================================================
    # compare data between scanned and system POs Start
    # ============================================================
    async def compare_sharepoint_scanned_and_system_pos(
        self,
        request=None,
        app=None,
        user_id: int = None,
        sharepoint_po_det_ids: list[int] = None,
        sp_repo=SharepointRepo,
        system_pos=None
    ):
        try:
            # ---------------- Resolve app context ---------------- #
            resolved_app = app or (request.app if request is not None else None)

            if resolved_app is None:
                logger.error("No app context provided to compare_sharepoint_scanned_and_system_pos")
                return {
                    "status": "error",
                    "message": "Failed to generate PO report",
                    "error": "No app context available"
                }
            
            # ---------------- Fetch scanned POs ---------------- #
            scanned_pos = await sp_repo.get_sharepoint_po_details_by_ids(sharepoint_po_det_ids)

            if not scanned_pos:
                return {
                    "status": "success",
                    "message": "No scanned POs found for comparison"
                }

            scanned_pos = [
                {k: self.make_json_safe(v) for k, v in po.items()}
                for po in scanned_pos
            ]

            # ---------------- Fetch system POs ---------------- #
            scanned_po_numbers = list({
                po["po_number"] for po in scanned_pos if po.get("po_number")
            })

            # -------------------- Fetch system POs -------------------- #
            # oldest_date = await sp_repo.get_sharepoint_oldest_report_date()

            # if oldest_date:
            #     system_pos = await MSSQLRepo.get_po_list(resolved_app, oldest_date)
            # else:
            #     system_pos = await MSSQLRepo.get_po_list_without_oldest_date(resolved_app)

            # # ---- imaginary PK for system_pos ---- #
            # for po in system_pos:
            #     po["system_po_id"] = self.make_stable_system_po_id(po)

            # system_pos = [
            #     {k: self.make_json_safe(v) for k, v in po.items()}
            #     for po in system_pos
            # ]

            #-----------reconcile OLD mismatches/missing BEFORE processing new POs----------
            # await self.reconcile_all_pos(
            #     user_id=user_id,
            #     system_pos=system_pos,
            #     sp_repo=sp_repo
            # )
            
            # ---------------- Create system PO lookup ---------------- #
            system_po_map = defaultdict(list)

            for po in system_pos:
                po_number = po.get("po_number")
                if not po_number:
                    continue
                key = self.normalize_po(po_number)
                system_po_map[key].append(po)

            matched_pairs = []
            missing_pos = []

            # ---------------- Match scanned with system ---------------- #
            for scanned in scanned_pos:

                scanned_po = scanned.get("po_number")
                normalized_po = self.normalize_po(scanned_po)

                candidates = system_po_map.get(normalized_po)

                if not candidates:
                    missing_pos.append(scanned)
                    continue

                # choose best candidate
                system = self.find_best_system_match(scanned, candidates)

                if not system:
                    missing_pos.append(scanned)
                    continue

                matched_pairs.append({
                    "sharepoint_po_det_id": scanned["sharepoint_po_det_id"],
                    "system_po_id": system["system_po_id"],
                    "scanned": {f: scanned.get(f) for f in self.FIELDS_TO_COMPARE},
                    "system": {f: system.get(f) for f in self.FIELDS_TO_COMPARE},
                    "raw_scanned": scanned
                })

            # ---------------- Prepare pairs for OpenAI ---------------- #
            pairs_for_llm = []

            for pair in matched_pairs:

                for field in self.FIELDS_TO_COMPARE:

                    scanned_val = pair["scanned"].get(field)
                    system_val = pair["system"].get(field)

                    if scanned_val in (None, "") or system_val in (None, ""):
                        continue

                    scanned_norm = self.normalize_value(scanned_val)
                    system_norm = self.normalize_value(system_val)

                    if scanned_norm == system_norm:
                        continue

                    pairs_for_llm.append({
                        "sharepoint_po_det_id": pair["sharepoint_po_det_id"],
                        "system_po_id": pair["system_po_id"],
                        "field": field,
                        "scanned_value": str(scanned_val),
                        "system_value": str(system_val)
                    })

            # ---------------- Call OpenAI ---------------- #
            mismatches = []

            if pairs_for_llm:
                mismatches = await self.llm_batch_compare(pairs_for_llm)

            mismatch_pairs = set()

            # ---------------- Insert mismatches ---------------- #
            for mm in mismatches:

                scanned_value = "" if mm["scanned_value"] is None else str(mm["scanned_value"])
                system_value = "" if mm["system_value"] is None else str(mm["system_value"])

                exists = await sp_repo.mismatch_exists(
                    user_id=user_id,
                    sharepoint_po_det_id=mm["sharepoint_po_det_id"],
                    system_po_id=mm["system_po_id"],
                    mismatch_attribute=mm["field"],
                    scanned_value=scanned_value,
                    system_value=system_value
                )

                if not exists:

                    await sp_repo.insert_mismatch(
                        sharepoint_po_det_id=mm["sharepoint_po_det_id"],
                        user_id=user_id,
                        system_po_id=mm["system_po_id"],
                        field=mm["field"],
                        system_value=system_value,
                        scanned_value=scanned_value,
                        comment=f"{mm['field']} mismatch"
                    )

                mismatch_pairs.add((mm["sharepoint_po_det_id"], mm["system_po_id"]))

            # ---------------- Insert missing fields ---------------- #
            for pair in matched_pairs:

                for field in FIELDS_TO_COMPARE:

                    scanned_val = pair["scanned"].get(field)
                    system_val = pair["system"].get(field)

                    if system_val not in (None, "") and scanned_val in (None, ""):

                        exists = await sp_repo.mismatch_exists(
                            user_id=user_id,
                            sharepoint_po_det_id=pair["sharepoint_po_det_id"],
                            system_po_id=pair["system_po_id"],
                            mismatch_attribute=field,
                            scanned_value="",
                            system_value=str(system_val)
                        )

                        if not exists:

                            await sp_repo.insert_mismatch(
                                sharepoint_po_det_id=pair["sharepoint_po_det_id"],
                                user_id=user_id,
                                system_po_id=pair["system_po_id"],
                                field=field,
                                system_value=str(system_val),
                                scanned_value="",
                                comment=f"{field} missing in scanned data"
                            )

                        mismatch_pairs.add((pair["sharepoint_po_det_id"], pair["system_po_id"]))

            # ---------------- Insert matched records ---------------- #
            for pair in matched_pairs:

                key = (pair["sharepoint_po_det_id"], pair["system_po_id"])

                if key in mismatch_pairs:
                    continue

                scanned = pair["raw_scanned"]

                await sp_repo.insert_sharepoint_matched_po(
                    sharepoint_po_det_id=pair["sharepoint_po_det_id"],
                    system_po_id=pair["system_po_id"],
                    sharepoint_file_id=scanned.get("sharepoint_file_id"),
                    user_id=user_id,
                    po_number=scanned.get("po_number"),
                    po_date=scanned.get("po_date"),
                    vendor_number=scanned.get("vendor_number"),
                    customer_name=scanned.get("customer_name"),
                    created_by="system"
                )

            # ---------------- Insert PO missing ---------------- #
            for po in missing_pos:

                exists = await sp_repo.po_missing_exists(
                    user_id=user_id,
                    sharepoint_po_det_id=po["sharepoint_po_det_id"],
                    system_po_id=None,
                    mismatch_attribute="po_missing",
                    scanned_value=po.get("po_number"),
                    system_value=""
                )

                if not exists:

                    await sp_repo.insert_po_missing(
                        sharepoint_po_det_id=po["sharepoint_po_det_id"],
                        user_id=user_id,
                        system_po_id=None,
                        attribute="po_missing",
                        system_value="",
                        scanned_value=po.get("po_number"),
                        comment="PO not found in system"
                    )

            return {
                "status": "success",
                "message": "PO comparison completed successfully"
            }

        except Exception as e:
            logger.exception(
                f"Error in compare_sharepoint_scanned_and_system_pos | user_id={user_id}"
            )

            return {
                "status": "error",
                "message": "Failed to generate PO report",
                "error": str(e)
            }
    # ============================================================
    # compare data between scanned and system POs End
    # ============================================================