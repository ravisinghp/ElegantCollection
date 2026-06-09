
from dotenv import load_dotenv
import os, base64
import httpx
from urllib.parse import urlencode
import json
from decimal import Decimal
from loguru import logger
import jwt
from datetime import date, datetime
from openai import OpenAI
import asyncio
from app.utils.image_ocr import extract_text_from_image_bytes
from app.db.repositories.sync_client_po_repo import MSSQLRepo
from collections import defaultdict

# Load the .env file
load_dotenv()
failed_url = os.getenv("failed_url")
success_url = os.getenv("success_url")

#----------------- outlook ------------------#
CLIENT_ID = os.getenv("CLIENT_ID")
CLIENT_SECRET = os.getenv("CLIENT_SECRET")
TENANT_ID = os.getenv("TENANT_ID")
REDIRECT_URI = os.getenv("REDIRECT_URI")
GRAPH_API = os.getenv("GRAPH_API")
JWTSECRET_KEY=os.getenv("JWTSECRET_KEY")
#---------------outlook end ----------------------#

#---------------Google---------------------------#
GOOGLE_CLIENT_ID = os.getenv("GOOGLE_CLIENT_ID")
GOOGLE_CLIENT_SECRET = os.getenv("GOOGLE_CLIENT_SECRET")
# TENANT_ID = os.getenv("TENANT_ID")
GOOGLE_REDIRECT_URI = os.getenv("GOOGLE_REDIRECT_URI")
GMAIL_API = os.getenv("GMAIL_API")
#---------------Google end-------------------------#

#---------------OpenAI Client------------------
openai_client = OpenAI(
    api_key=os.getenv("OPENAI_API_KEY")
)
#----------OpenAI Client end ------------------

from datetime import datetime,timedelta
import re
from typing import List, Dict, Any, Optional
import html
import io
import aiohttp, json
from typing import List, Dict, Any
from pptx import Presentation  # for PPTX support
from fastapi.responses import RedirectResponse
try:
    # Optional import to avoid circular issues in other contexts
    from app.db.repositories.mails import MailsRepository
except Exception:
    MailsRepository = None  # type: ignore

try:
    import PyPDF2  # type: ignore
except Exception as e:
    print("PyPDF2 import failed:", e)
    PyPDF2 = None  # type: ignore

try:
    import docx  # python-docx
except Exception:
    docx = None  # type: ignore


#------------------Block domain------------------
BLOCKED_DOMAINS = {
    domain.strip().lower()
    for domain in os.getenv("BLOCKED_MAIL_DOMAINS", "").split(",")
    if domain.strip()
}
#------------------Block domain------------------


def get_auth_url(provider: str, user_id: int):

    state = jwt.encode(
        {"user_id": user_id},
        JWTSECRET_KEY,
        algorithm="HS256"
    )
    
    if isinstance(state, bytes):
        state = state.decode("utf-8")
      
    if provider == "outlook":
        return (
            f"https://login.microsoftonline.com/{TENANT_ID}/oauth2/v2.0/authorize?"
            f"client_id={CLIENT_ID}"
            f"&response_type=code"
            f"&redirect_uri={REDIRECT_URI}"
            f"&response_mode=query"
            f"&scope=offline_access%20Mail.Read%20User.Read"
            f"&state={state}"
            f"&prompt=login"
        )
    elif provider == "google":
        params = {
            "client_id": GOOGLE_CLIENT_ID,
            "redirect_uri": GOOGLE_REDIRECT_URI,
            "response_type": "code",
            "scope": "https://www.googleapis.com/auth/gmail.readonly https://www.googleapis.com/auth/calendar.readonly https://www.googleapis.com/auth/userinfo.email openid",
            "access_type": "offline",
            "prompt": "consent"
        }
        return f"https://accounts.google.com/o/oauth2/v2/auth?{urlencode(params)}"
    else:
        raise ValueError("Invalid provider")


# ------------------- token exchange -----------------------
async def exchange_code_for_token(code: str):
    url = f"https://login.microsoftonline.com/{TENANT_ID}/oauth2/v2.0/token"
    data = {
        'client_id': CLIENT_ID,
        'scope': 'offline_access Mail.Read',  # must match authorization URL
        'code': code,
        'redirect_uri': REDIRECT_URI,
        'grant_type': 'authorization_code',
        'client_secret': CLIENT_SECRET
    }
    headers = {"Content-Type": "application/x-www-form-urlencoded"}

    async with httpx.AsyncClient(timeout=httpx.Timeout(10.0, read=30.0)) as client:
        response = await client.post(url, data=data, headers=headers)
        token_json = response.json()

    if response.status_code != 200 or "access_token" not in token_json:
        return {"error": "Token exchange failed", "details": token_json}

    return {
        "access_token": token_json.get("access_token"),
        "refresh_token": token_json.get("refresh_token"),
        "expires_in": token_json.get("expires_in"),
        "url": f"{success_url}?mail_token={token_json.get('access_token')}"
    }
    
    
# this code is used to generate token for gmail(Google api)
async def exchange_code_for_token_for_gmail(code: str) -> dict:
        """Exchange auth code for access/refresh tokens"""
        token_url = "https://oauth2.googleapis.com/token"
        data = {
            "code": code,
            "client_id": GOOGLE_CLIENT_ID,
            "client_secret": GOOGLE_CLIENT_SECRET,
            "redirect_uri": GOOGLE_REDIRECT_URI,
            "grant_type": "authorization_code",
        }

        async with httpx.AsyncClient() as client:
            headers = {"Content-Type": "application/x-www-form-urlencoded"}
            resp = await client.post(token_url, data=data, headers=headers)
            return resp.json()


# this code is used to fetch email from google api
async def get_user_email(access_token: str) -> str | None:
        """Fetch user email from Google API"""
        async with httpx.AsyncClient() as client:
            resp = await client.get(
                "https://www.googleapis.com/oauth2/v2/userinfo",
                headers={"Authorization": f"Bearer {access_token}"}
            )
            data = resp.json()
            return data.get("email")

###------------------This code is used to fetch folder names upto 200 folders------------------
async def fetch_all_folders(access_token: str) -> List[Dict[str, Any]]:
    headers = {"Authorization": f"Bearer {access_token}"}
    timeout = httpx.Timeout(connect=10.0, read=60.0, write=30.0, pool=30.0)

    folder_list = []
    url = f"{GRAPH_API}/me/mailFolders?$top=200&$expand=childFolders"

    async with httpx.AsyncClient(timeout=timeout) as client:
        while url:
            resp = await client.get(url, headers=headers)
            resp.raise_for_status()
            data = resp.json()

            for folder in data.get("value", []):
                folder_list.append({
                    "id": folder.get("id"),
                    "name": folder.get("displayName")
                })

            # If Graph gives a nextLink, continue paging
            url = data.get("@odata.nextLink")

    return folder_list

# mansi-------------------------------------------------------------

async def refresh_outlook_access_token(refresh_token: str) -> dict:
    url = "https://login.microsoftonline.com/common/oauth2/v2.0/token"

    data = {
        "client_id": CLIENT_ID,
       # "client_secret": CLIENT_SECRET,
        "grant_type": "refresh_token",
        "refresh_token": refresh_token,
        #"scope": "https://graph.microsoft.com/.default offline_access",
        # "scope": "offline_access User.Read Mail.Read Mail.ReadWrite"
    }

    timeout = httpx.Timeout(connect=10.0, read=30.0, write=30.0, pool=30.0)

    async with httpx.AsyncClient(timeout=timeout) as client:
        resp = await client.post(url, data=data)
        print(f"the response of data : {data}")
        resp.raise_for_status()
        return resp.json()


# mansi --------------------------------------------------------------------------------------------------------
async def get_valid_outlook_token(
    user_id: int,
    repo: MailsRepository,
) -> str | None:
    token = await repo.get_outlook_token(user_id)
    if not token:
        return None

    if token.token_expiry > datetime.utcnow() + timedelta(minutes=2):
        return token.access_token

    try:
        new_tokens = await refresh_outlook_access_token(token.refresh_token)
        print(f"the response of updated refresh token : {token}")
        print(f"the response of updated refresh new token : {new_tokens}")
    except Exception as e:
        print("Exception:",e)
        # refresh failed, maybe refresh token revoked
        return None

    new_expiry = datetime.utcnow() + timedelta(seconds=int(new_tokens["expires_in"]))

    await repo.update_outlook_token(
        user_id=user_id,
        access_token=new_tokens["access_token"],
        refresh_token=new_tokens.get("refresh_token", token.refresh_token),
        token_expiry=new_expiry,
    )

    return new_tokens["access_token"]


# ------------------Email + Attachment Fetching + LLM logic start ------------------ #
def strip_html_to_text(html_content: Optional[str]) -> str:
    if not html_content:
        return ""
    text = re.sub(r"<[^>]+>", " ", html_content)
    text = html.unescape(text)
    return re.sub(r"\s+", " ", text).strip()


def clean_email_body(body_text: str) -> str:
    if not body_text:
        return ""
    lines = body_text.split("\n")
    cleaned_lines: List[str] = []
    skip_line_starts = [
        r"^From:", r"^To:", r"^Cc:", r"^CC:", r"^BCC:", r"^Bcc:", r"^Sent:",
        r"^Subject:", r"^Reply-To:", r"^Message-ID:", r"^X-.*?:",
        r"^Content-Type:", r"^Content-Transfer-Encoding:", r"^MIME-Version:",
        r"^Return-Path:", r"^Delivered-To:", r"^Received:",
        r"^On .* wrote:", r"^-----Original Message-----",
        r"^Microsoft Teams$", r"^Need help\?$", r"^Join the meeting now$",
        r"^Meeting ID:", r"^Passcode:", r"^For organisers:",
        r"^Meeting options$", r"^_{6,}$",
    ]
    # REMOVED bare "^Date:" — it kills "PO Date: ..." lines.
    # Email header Date is already handled by the mail graph API fields.
    url_re = re.compile(r"^(https?://|www\.)", re.IGNORECASE)
    skip_res = [re.compile(pat, re.IGNORECASE) for pat in skip_line_starts]
    for raw_line in lines:
        line = raw_line.strip()
        if not line:
            continue
        if url_re.match(line):
            continue
        if any(rx.match(line) for rx in skip_res):
            continue
        cleaned_lines.append(line)
    cleaned_text = " ".join(cleaned_lines)
    return re.sub(r"\s+", " ", cleaned_text).strip() or body_text.strip()



def iso_to_date(iso_dt: Optional[str]) -> Optional[str]:
    if not iso_dt:
        return None
    return iso_dt[:10]


def collect_addresses_from_message(msg: Dict[str, Any], key: str) -> Optional[str]:
    out: List[str] = []
    for rec in msg.get(key, []) or []:
        address = (rec.get('emailAddress') or {}).get('address')
        if address:
            out.append(address)
    return ",".join(out) if out else None


def compute_file_hash(content: bytes) -> str:
    import hashlib
    return hashlib.sha256(content).hexdigest()


# ------------------- normalize text ------------------- #
def normalize_text(text: str) -> str:
    if not text:
        return ""
    text = text.replace("\xa0", " ")
    text = text.replace("\u200b", " ")
    text = text.replace("\r", "\n")
    # REMOVED: re.sub(r"[=]", ":", text) — corrupts style codes and URLs
    text = re.sub(r"\s+", " ", text)
    return text.strip()


def extract_po_fields_regex(text: str) -> dict:
    if not text or len(text) < 30:
        return EMPTY_PO

    text = normalize_text(text)
    out = EMPTY_PO.copy()

    for field, patterns in PO_REGEX_PATTERNS.items():
        for pat in patterns:
            match = re.search(pat, text, re.IGNORECASE)
            if match:
                # check if the match has at least 1 capturing group
                if match.lastindex and match.lastindex >= 1:
                    out[field] = match.group(1).strip()
                else:
                    # fallback: if no group, use full match
                    out[field] = match.group(0).strip()
                break

    return out if any(out.values()) else EMPTY_PO


DB_FIELD_LIMITS = {
    "customer_name": 255,
    "vendor_number": 100,
    "po_number": 100,
    "gold_karat": 100,
    "ec_style_number": 100,
    "customer_style_number": 100,
    "color": 50,
    "quantity": 100,
    "gold_lock": 100,
}
 
def trim_to_db_limits(data: dict) -> dict:
    trimmed = {}
 
    for field, value in data.items():
        if not value:
            trimmed[field] = value
            continue
 
        if field in DB_FIELD_LIMITS:
            max_len = DB_FIELD_LIMITS[field]
            trimmed[field] = value[:max_len]
        else:
            trimmed[field] = value
 
    return trimmed
  

MANDATORY_FIELDS = ["po_number", "customer_name", "vendor_number", "po_date", "delivery_date", "quantity", "cancel_date", "gold_karat", "ec_style_number", "customer_style_number", "color", "description", "gold_lock"]

async def extract_po_fields(text: str) -> dict:
    regex_data_response = extract_po_fields_regex(text)
    regex_data = trim_to_db_limits(regex_data_response)
    logger.info(f"regex data: {regex_data}")

    if all(regex_data.get(f) for f in MANDATORY_FIELDS) and len(text.strip()) >= 50:
        return regex_data

    if len(text.strip()) < 50:
        logger.info("Skipping LLM — text too short")
        return EMPTY_PO

    llm_data = await extract_po_fields_from_llm(text)
    logger.info(f"LLM data: {llm_data}")

    # Merge: LLM fills only the fields regex missed — never blanks a regex value
    final = llm_data
    # final = regex_data.copy()
    # for k, v in llm_data.items():
    #     if v not in (None, "", "null", "N/A"):
    #         if not final.get(k):          # only fill empty slots
    #             final[k] = v

    return final if any(final.values()) else EMPTY_PO


#--------------------Regex-----------------------------
PO_REGEX_PATTERNS = {

    # ---------------- PO NUMBER ----------------
    "po_number": [
        # ----- OLD WORKING -----
        r"(?:po_number|po_no)\s*:\s*(PO[\w\-_/]+)",
        r"(?:po\s*number|po\s*no|po#|po\s*#|p\.o\.|purchase\s*order|po)\s*[:\-]?\s*(PO[\w\-_/]+)",
        r"\b(PO[\s\-_:]*[0-9]{1,}[A-Z0-9\/_.\-]*)",
        r"(?:po\s*number|po\s*no|po#|p\.o\.|purchase\s*order)\s*[:\-]?\s*(PO[\- ]?[A-Z0-9\/_.\-]+)",
        # ----- NEW FORMATS -----
        r"(?:po\s*number|po\s*no|po#|p\.o\.|purchase\s*order)\s*[:\-]?\s*([A-Z]{1,5}-\d{4,}-\d+)",
        r"(?:po_number|po_no|po\s*number|po#|p\.o\.)\s*[:#]?\s*\n?\s*([A-Z0-9\-_/]+)",
    ],

    # ---------------- CUSTOMER NAME ----------------
    "customer_name": [
        r"(?i)\bcustomer\s*(?:name)?\b\s*\|\s*([A-Za-z0-9&.,\-/ ]{2,})",
        r"(?i)\bcustomer\s*(?:name)?\s*[:\-]\s*\n\s*([A-Za-z0-9&.,\-/ ]{2,})",
        r"(?i)\bcustomer\s*(?:name)?\s*[:\-]\s*([A-Za-z0-9&.,\-/ ]{2,})",
        r"(?i)\bship\s*to\s*:\s*\n\s*([A-Za-z0-9&.,\-/ ]{2,})",
        r"(?i)(?:ship\s*to|bill\s*to|deliver\s*to)\s*[:\-]\s*([A-Za-z0-9&.,\-/ ]{2,})",
        r"(?i)(?:customer|buyer|client)\s*[:\-]?\s*([A-Za-z0-9&.,\-/ ]+?)(?=\s+(?:vendor|po|date|delivery|qty|quantity|gold|color|description)\b|$)",
    ],

    # ---------------- VENDOR NUMBER ----------------
    "vendor_number": [
        r"Vendor\s*ID[\s\S]{0,80}\b(V\d{4,})\b",
        # Vendor Number / ID on SAME LINE
        r"(?:vendor[_\s]*(?:number|no|id)|supplier[_\s]*(?:number|no|code))\s*[:\-#]?\s*([A-Za-z0-9\-_./]+)",

        # Vendor ID on NEXT LINE 
        r"(?:vendor\s*id|vendor\s*number)\s*[:\-]?\s*\n\s*([A-Za-z0-9\-_./]+)",

        # Simple "Vendor: XYZ"
        r"\bvendor\b\s*[:\-]?\s*([A-Za-z0-9\-_./]+)",

        # V-ID / VNo formats
        r"\b(?:Vendor\s*ID|VNo|V-ID)\s*[:#\-\s]?\s*([A-Za-z0-9\-_./]+)",

        # Fallback supplier code
        r"(?:supplier\s*(?:no|number|code))\s*[:\-]?\s*([A-Za-z0-9\-_./]+)",
    ],                                                              

    # ---------------- PO DATE ----------------
    "po_date": [
        # ----- OLD WORKING -----
        r"(?:po\s*date|order\s*date|date)\s*[:\-]?\s*(\d{4}-\d{1,2}-\d{1,2})",
        r"po_date\s*:\s*(\d{4}-\d{2}-\d{2})",
        r"date\s*[:\-]?\s*(\d{4}-\d{2}-\d{2})",
        # ----- NEW FORMATS -----
        r"(?:Purchase\s+Order\s+Date|P\.O\.\s*Date)\s*[:\-]?\s*(\d{1,2}/\d{1,2}/\d{2})",
        r"P\.?\s*O\.?\s*Date\s*[:\-]?\s*(\d{1,2}-[A-Za-z]{3}-\d{4})",
        r"(?:purchase\s*order\s*date)\s*\n\s*(\d{1,2}/\d{1,2}/\d{2})",
        r"(?:purchase\s*order\s*date)\s*[:\-]?\s*(\d{1,2}-[A-Za-z]{3}-\d{4})",
        r"(?:po\s*date|order\s*date|date)\s*[:\-]?\s*([0-9]{1,2}-[A-Za-z]{3}-[0-9]{4})",
        r"P\.O\.\s+Date\s*:\s*(\d{2}-[A-Za-z]{3}-\d{4})",
        r"P\.O\.\s*Date\s*[:\-]?\s*(\d{1,2}-[A-Za-z]{3}-\d{4})",
        r"(?:po\s*date|order\s*date|date)\s*[:\-]?\s*(\d{4}-\d{1,2}-\d{1,2})",
        r"\b(\d{1,2}/\d{1,2}/\d{2})\b",
    ],

    #---------------- DELIVERY DATE ----------------
    "delivery_date": [
        # PDF table FIRST
        r"\b(?:DELIVERY\s*DATE|DUE\s*DATE|delivery_date)\b[\s\S]{0,100}?[:\s]*([\dA-Za-z/.-]{4,20})",

        # Inline / email
        r"(?:delivery\s*date|expected\s*delivery|delivery_date|due\s*date)\s*[:\-]?\s*(\d{4}-\d{2}-\d{2})",
        r"(?:delivery\s*date|expected\s*delivery|delivery_date|due\s*date)\s*[:\-]?\s*(\d{2}-[A-Za-z]{3}-\d{4})",
    ], 
 
    # ---------------- CANCEL DATE ----------------
    "cancel_date": [
        r"(?:cancel\s*date|cancellation\s*date)\s*[:\-]?\s*(\d{4}-\d{1,2}-\d{1,2})",
        r"cancel_date\s*:\s*(\d{4}-\d{2}-\d{2})",
    ],
 
    # ---------------- EC STYLE NUMBER ----------------
   "ec_style_number": [
    r"(?:ec\s*style\s*number|ec\s*style|ec\s*no)\s*[:\-]?\s*([A-Z0-9.\-]+)",
    r"(?:ec_style_number|ec_style_no)\s*:\s*([A-Z0-9.\-]+)",
    ],
 
    # ---------------- CUSTOMER STYLE NUMBER ----------------
    "customer_style_number": [
        r"(?:customer\s*style\s*number|customer\s*style|cust\s*style)\s*[:\-]?\s*([A-Z0-9\-]+)",
        r"(?:customer_style_number|customer_style_no)\s*:\s*([A-Z0-9\-]+)",
    ],
 
    # ---------------- QUANTITY ----------------
    "quantity": [
        # Table style (PDF / same line)
        r"\bQUANTITY\b[\s\S]{0,100}\n\s*([A-Za-z0-9 ,\-–\.]{10,})",

        # Vertical layout (Quantity\n1)
        r"(?i)\bquantity\b\s*[\r\n]+\s*(\d+)\b",

        # EA / PCS rows
        r"\n\s*(\d+)\s+(?:EA|PCS|PC)\b",

        # Inline
        r"(?i)(?:qty|quantity|pcs|pieces)\s*[:\-]?\s*(\d+)",

        # Fallback (keep last!)
        r"(?i)\b(\d+)\s*(?:pcs|pieces|nos)\b",
    ],
 
    # ---------------- GOLD KARAT ----------------
    "gold_karat": [
        r"(?i)\bMETAL\b[\s\S]{0,60}?\b([^\s|,\n\r]+)",
        r"(?i)(?:gold\s*karat|gold_karat|gold\s*purity|karat|metal|kt)\s*[:\-]?\s*([^\n\r]+?)(?=\s+(?:ec\s*style|customer\s*style|quantity|color|description|gold\s*lock)\s*[:\-]|$)",
        r"(?i)(?:gold\s*karat|gold_karat|karat|metal)\s*[\r\n]+\s*([^\r\n]+)",
    ],              
 
    # ---------------- COLOR ----------------
   "color": [
        r"(?:color|colour)\s*[:\-]?\s*([A-Za-z0-9\s\-]+)",
        r"(?:color|colour)\s*[:\-]?\s*([A-Za-z0-9\s\-]+?)(?=\s+(?:quantity|gold|karat|description|remarks|details)\s*:|$)",
        r"color\s*[:\-]?\s*([A-Za-z0-9\s\-]+)"
    ],

    #---------------- DESCRIPTION ----------------
    "description": [
        # PDF table
        r"\bDESCRIPTION\b[\s\S]{0,100}\n\s*([A-Za-z0-9 ,\-–\.]{10,})",

        # Structured lines
        r"(?:item\s*description|description)\s*[:\-]?\s*([A-Za-z][A-Za-z\s\-–]+)",

        # Item row patterns
        r"[A-Z0-9\-]+\s+\d+KW\s+([A-Za-z ].*?SIZE:\s*[0-9.]+)",
    ],

    # ---------------- GOLD LOCK ----------------
    "gold_lock": [
        r"(?:gold\s*lock|gold_lock|gold\s*locking|lock\s*value|lock\s*percentage|metal\s*lock|lock)\s*[:\-]?\s*([0-9]+(?:\.[0-9]+)?)",
        r"(?:gold\s*lock|gold_lock|gold\s*locking|lock\s*value|lock\s*percentage|metal\s*lock|lock)\s*[:\-]?\s*[\n\r]+?\s*([0-9]+(?:\.[0-9]+)?)",
        r"(?:gold\s*lock|gold_lock|gold\s*locking|lock\s*value|lock\s*percentage|metal\s*lock|lock)\s{2,}([0-9]+(?:\.[0-9]+)?)",
        r"(?:goldlock|gold_lock|lock)\s*([0-9]+(?:\.[0-9]+)?)"
    ],

}


EMPTY_PO = {
    "po_number": None,
    "customer_name": None,
    "vendor_number": None,
    "po_date": None,
    "delivery_date": None,
    "cancel_date": None,
    "gold_karat": None,
    "ec_style_number": None,
    "customer_style_number": None,
    "color": None,
    "quantity": None,
    "description": None,
    "gold_lock": None,
}

KEYWORD_REGEX_MAP = {
    "po_number": [
        r"\bpo\s*(no|number|#|id)\b",
        r"\bp\.o\.\s*(no|number|#)?\b",
        r"\bpurchase\s*order\s*(no|number)?\b",
        r"\b(po)\s*(?:no|number|#|id)?\s*\d+\b",
        r"\b(po)\d+\b",
        r"\b(p\.o\.)\s*(?:no|number|#)?\s*\d*\b",
        r"\b(purchase\s*order)\s*(?:no|number)?\s*\d*\b",
        r"\b(order)\s*(?:no|number|#|id)?\s*\d+\b",
       r"\b(order|orders|ordered|reorder|re-?order|ordering|purchase|purchases|po|attached|attachment|attachments|enclosed)\b"
    ],

    "customer_name": [
        r"\bcustomer\s*name\b",
        r"\bbuyer\b",
        r"\bclient\b",
        r"\bparty\s*name\b"
    ],

    "vendor_number": [
        r"\bvendor\s*(no|number|code)\b",
        r"\bsupplier\s*(no|number|code)\b"
    ],

    "po_date": [
        r"\bpo\s*date\b",
        r"\border\s*date\b",
        r"\bdate\s*of\s*order\b"
    ],

    "delivery_date": [
        r"\bdelivery\s*date\b",
        r"\bexpected\s*delivery\b",
        r"\bdispatch\s*date\b"
    ],

    "cancel_date": [
        r"\bcancel\s*date\b",
        r"\bexpiry\s*date\b"
    ],

    "gold_karat": [
        r"\b\d{2}\s*K\b",
        r"\b\d{2}\s*KT\b",
        r"\bAG\s*\d{3}\b",
        r"\bgold\s*karat\b",
        r"\bpurity\b"
    ],

    "ec_style_number": [
        r"\bec\s*style\s*(no|number)\b",
        r"\bec\s*style\b"
    ],

    "customer_style_number": [
        r"\bcustomer\s*style\s*(no|number)\b",
        r"\bdesign\s*(no|number)\b"
    ],

    "color": [
        r"\b(yellow|white|rose)\s*gold\b",
        r"\bcolor\b"
    ],

    "quantity": [
        r"\bqty\b",
        r"\bquantity\b",
        r"\bpcs\b",
        r"\bnos\b",
        r"\bpieces\b"
    ],

    "description": [
        r"\bdescription\b",
        r"\bitem\s*details\b",
        r"\bproduct\s*details\b"
    ],

    "gold_lock": [
        r"\bgold\s*lock\b",
        r"\blocked\s*gold\b",
        r"\block\s*status\b"
    ],
}

async def detect_keywords(text: str, db_keywords: list[str]):
    if not text or not text.strip():
        return [], None

    text_l = text.lower()
    detected = set()

    # ---------------- 1. DB keywords (exact-ish) ----------------
    for kw in db_keywords:
        if normalize_keyword(kw) in text_l:
            detected.add(kw)

    # ---------------- 2. Jewellery-aware regex detection ----------------
    for field, patterns in KEYWORD_REGEX_MAP.items():
        for pat in patterns:
            if re.search(pat, text_l, re.IGNORECASE):
                detected.add(field)
                break

    if detected:
        return sorted(detected), "REGEX_MATCH"

    return [], None

def normalize_keyword(k: str) -> str:
    return re.sub(r"\s+", " ", k.strip().lower())


# ------------------- Extraction ------------------- #
PO_FIELD_NAMES = [
    "po_number",
    "customer_name",
    "vendor_number",
    "po_date",
    "delivery_date",
    "cancel_date",
    "ec_style_number",
    "customer_style_number",
    "quantity",
    "gold_karat",
    "color",
    "description",
    "gold_lock"
]

EMPTY_PO = {field: None for field in PO_FIELD_NAMES}


async def extract_po_fields_from_llm(text: str) -> dict:
    field_list = json.dumps(PO_FIELD_NAMES, indent=2)
    prompt = f"""
You are a precision extraction engine for Jewelry Purchase Orders.

Your job is to extract structured data from raw text parsed from files
(PDFs, Excel sheets, Word docs, images via OCR, or email bodies).

===========================
GLOBAL RULES (NON-NEGOTIABLE)
===========================
- Extract ONLY values that are EXPLICITLY present with a clear label.
- NEVER infer, guess, or derive values from context.
- NEVER merge two fields into one value.
- NEVER include field label names inside field values.
- If a value is absent, ambiguous, or only implied return null.
- Preserve original formatting (casing, spacing, units) in extracted values.
- Do not normalize, rewrite, or clean values.

===========================
SOURCE-SPECIFIC PARSING RULES
===========================

[EXCEL / CSV TEXT]
- Text arrives as row-by-row dumps separated by pipe characters.
- First row is typically the header row — match column names to fields.
- If multiple data rows exist extract the FIRST valid row only (unless items list).
- Ignore "Sheet: name" prefix lines.

[PDF / OCR TEXT]
- Labels and values may be on separate lines.
- For key:value pairs the value is text immediately AFTER the colon on the same line,
  or the first non-empty line below the label.

[EMAIL BODY TEXT]
- Extract only from the FIRST (top) email block.
- Ignore quoted previous messages and signatures.
- Labels may be inline e.g. "PO#: 12345, Customer: ABC Corp" — split correctly.

[IMAGE / OCR TEXT]
- Apply light OCR error tolerance: "P0 Number" likely means "PO Number".
- Only extract if you are more than 90 percent certain of the value.

===========================
FIELD-BY-FIELD RULES
===========================

PO NUMBER:
- Labels: "PO Number", "PO No", "P.O.", "PO#", "PO #", "Purchase Order Number", "Order No"
- Value must look like an alphanumeric PO identifier e.g. "PO-2024-001", "12345A", "EC\SO\20\SV1\3\7".
- If multiple PO numbers appear return the FIRST one.

CUSTOMER NAME:
════════════════════════════════════════════════
CORE LOGIC (works for ANY document format):
════════════════════════════════════════════════

This PO was sent BY a customer TO Elegant Collection (the vendor).
Your job: find WHO SENT this PO.

ELEGANT COLLECTION is ALWAYS the vendor — NEVER return it as customer name.
Any variation of "Elegant Collection" must be rejected as customer name.

════════════════════════════════════════════════
SIGNALS THAT IDENTIFY THE CUSTOMER (issuer):
════════════════════════════════════════════════

STRONG signals — high confidence:
- Company name/logo text in the FIRST 3 lines of the document
- Text near labels: "From:", "Issued By:", "Purchaser:", "Buyer:", "Prepared By:"
- Company name in the document FOOTER (repeated at bottom of every page)
- Company name in legal/terms text:
  e.g. "XYZ Corp reserves the right to cancel..."
  e.g. "transactions with XYZ Corp and/or its subsidiaries"
  e.g. "governed by XYZ Corp"
- Email domain in contact info:
  e.g. "buyer@davidyurman.com" → David Yurman
  e.g. "BackOffice@Lsdco.com" → Leo Schachter

MEDIUM signals — use if strong signals absent:
- Company name on the BILL TO address block
  (the company billing = usually the customer)
- Letterhead address that does NOT match Elegant Collection's address
  (Plot 56A, SEEPZ, Mumbai = Elegant Collection — reject this)
- Company name repeated multiple times in the document

WEAK signals — last resort only:
- Company name in "Ship To:" block
- Company name embedded in item descriptions

════════════════════════════════════════════════
SIGNALS TO ALWAYS REJECT AS CUSTOMER NAME:
════════════════════════════════════════════════
- "ELEGANT COLLECTION" or any variation → this is the vendor
- "Order To:" field value → this is the vendor 
- "Ship To:" field value → this is a shipping destination
- "Vendor:" field value → this is the vendor
- Any street address, city, state, ZIP, country
- Phone numbers, fax numbers, email addresses
- Generic words: "purchaser", "buyer", "vendor", "supplier" alone

════════════════════════════════════════════════
FORMAT-SPECIFIC HINTS:
════════════════════════════════════════════════
- If document starts with a company name on line 1 or 2 → very likely the customer
- If "PURCHASE ORDER" header exists → issuing company is above or below it
- If "Jewelry Vendor P.O." or similar title → the company in the header/footer issued it
- If only one non-Elegant-Collection company name exists in the whole document → that is the customer
- Vend.PO# / Vendor PO # → the number after this label is the PO number, company above it is customer

════════════════════════════════════════════════
RETURN RULES:
════════════════════════════════════════════════
- Return the SHORTEST clean company name found
  e.g. "LEO SCHACHTER NEW YORK" not "LEO SCHACHTER (NEW YORK) INC. 50 West 47th..."
  e.g. "DAVID YURMAN" not "David Yurman Enterprises LLC"
- Strip legal suffixes if the short name is clear: LLC, INC, LTD, CORP
- Return null if genuinely cannot determine — do NOT guess

VENDOR NUMBER:
- Labels: "Vendor", "Vendor No", "Vendor Number", "Supplier Code", "Vendor ID"
- Value is typically numeric or alphanumeric e.g. "V-1042", "5583".

PO DATE:
- Labels: "PO Date", "Order Date", "Purchase Order Date", "Date"
- Accept bare "Date" only if near PO Number or document title.
- Return date exactly as written.

DELIVERY DATE:
- Labels: "Delivery Date", "Due Date", "Ship By", "Required By", "Expected Delivery"
- Return null if only per-item delivery dates exist.

CANCEL DATE:
- Labels: "Cancel Date", "Cancellation Date", "Cancel By", "Void After"

GOLD KARAT:
- Labels: "Gold Karat", "Karat", "Metal", "Gold Purity", "KT", "K"
- Valid examples: "14KW", "18K", "22K Yellow", "18 KT", "14K Rose Gold"
- Label MUST be a standalone field header not embedded in a description.

EC STYLE NUMBER:
- Labels: "EC Style", "EC Style No", "EC Style Number", "Style#", "IDD Style #"
- Alphanumeric style code. Do NOT extract customer-side style numbers.

CUSTOMER STYLE NUMBER:
- Labels: "Customer Style", "Cust Style", "Your Style", "Vendor Style #"
- This is the buyer's own reference number distinct from EC Style.

COLOR:
- Labels: "Color", "Colour", "Metal Color"
- Do NOT extract color embedded in description sentences.

DESCRIPTION:
- Labels: "Description", "Item Description", "Particulars", "Goods Description"
- Extract the FULL value exactly as written.

QUANTITY:
- Labels: "Quantity", "Qty", "Pieces", "Pcs", "Ord#"
- Value must be numeric.
- Do NOT sum multiple item quantities — return null if only per-item quantities exist.

GOLD LOCK:
- Labels: "Gold Lock", "Lock Value", "Lock %", "Metal Lock"
- Value may be a percentage or numeric.

===========================
MULTI-ROW ITEM DETECTION
===========================
If the text contains a TABLE with multiple product/item rows return:
  "has_line_items": true
Still extract all header-level fields. Return null for per-item fields
(Delivery Date, Gold Karat, Description, Quantity, Color) in the main record.

===========================
OUTPUT FORMAT
===========================
Return a single valid JSON object with EXACTLY these keys:
{field_list}
Plus one optional key: "has_line_items" (boolean, default false).

Rules:
- All string values: preserve original casing and formatting.
- Missing or ambiguous values: null (not empty string, not "N/A").
- No explanation, preamble, or markdown — JSON only.

===========================
TEXT TO EXTRACT FROM:
===========================
{text}
"""
    try:
        resp = openai_client.chat.completions.create(
            model="gpt-4o-mini",
            temperature=0,
            messages=[{"role": "user", "content": prompt}]
        )
        raw = resp.choices[0].message.content.strip()
        raw = re.sub(r"^```(?:json)?|```$", "", raw, flags=re.MULTILINE | re.IGNORECASE).strip()
        match = re.search(r"\{.*\}", raw, re.DOTALL)
        if not match:
            return EMPTY_PO
        data = json.loads(match.group())
        result = {}
        for k in PO_FIELD_NAMES:
            v = data.get(k)
            result[k] = v if v not in ("", None, "null", "N/A") else None
        return result
    except Exception as e:
        logger.error(f"LLM extraction failed: {e}")
        return EMPTY_PO

from datetime import datetime
from typing import Optional

def normalize_po_date_ddmmyyyy(date_str: Optional[str]) -> Optional[str]:
    """
    Converts LLM or regex date output to YYYY-MM-DD string.
    Returns None if parsing fails.
    """
    if not date_str:
        return None

    date_str = date_str.strip()
    
    # ── pre-clean
    date_str = re.sub(r'(\d+)(st|nd|rd|th)\b', r'\1', date_str, flags=re.IGNORECASE)
    date_str = re.sub(r'[ \t]+', ' ', date_str).strip()

    date_formats = [
        "%Y-%m-%d",          # 2025-07-11
        "%Y/%m/%d",          # 2025/07/11
        "%Y.%m.%d",          # 2025.07.11
        "%Y %m %d",          # 2025 07 11

        "%d-%m-%Y",          # 11-07-2025
        "%d/%m/%Y",          # 11/07/2025
        "%d.%m.%Y",          # 11.07.2025
        "%d %m %Y",          # 11 07 2025

        "%m-%d-%Y",          # 07-11-2025
        "%m/%d/%Y",          # 07/11/2025
        "%m.%d.%Y",          # 07.11.2025
        "%m %d %Y",          # 07 11 2025

        "%y-%m-%d",          # 25-07-11
        "%y/%m/%d",          # 25/07/11
        "%y.%m.%d",          # 25.07.11

        "%d-%m-%y",          # 11-07-25
        "%d/%m/%y",          # 11/07/25
        "%d.%m.%y",          # 11.07.25
        "%d %m %y",          # 11 07 25

        "%m-%d-%y",          # 07-11-25
        "%m/%d/%y",          # 07/11/25
        "%m.%d.%y",          # 07.11.25
        "%m %d %y",          # 07 11 25

        "%d-%b-%Y",          # 11-Jul-2025  
        "%d/%b/%Y",          # 11/Jul/2025
        "%d.%b.%Y",          # 11.Jul.2025
        "%d %b %Y",          # 11 Jul 2025

        "%b-%d-%Y",          # Jul-11-2025
        "%b/%d/%Y",          # Jul/11/2025
        "%b %d %Y",          # Jul 11 2025
        "%b %d, %Y",         # Jul 11, 2025
        "%b-%d-%Y",          # APR-19-2026

        "%Y-%b-%d",          # 2025-Jul-11
        "%Y/%b/%d",          # 2025/Jul/11

        "%d-%b-%y",          # 19-APR-26  
        "%d/%b/%y",          # 19/APR/26
        "%d.%b.%y",          # 19.APR.26
        "%d %b %y",          # 19 APR 26

        "%b-%d-%y",          # APR-19-26
        "%b/%d/%y",          # APR/19/26
        "%b %d %y",          # APR 19 26
        "%b %d, %y",         # APR 19, 26

        "%d-%B-%Y",          # 11-July-2025
        "%d/%B/%Y",          # 11/July/2025
        "%d.%B.%Y",          # 11.July.2025
        "%d %B %Y",          # 11 July 2025

        "%B-%d-%Y",          # July-11-2025
        "%B/%d/%Y",          # July/11/2025
        "%B %d %Y",          # July 11 2025
        "%B %d, %Y",         # July 11, 2025

        "%Y-%B-%d",          # 2025-July-11
        "%Y/%B/%d",          # 2025/July/11

        "%d-%B-%y",          # 19-APRIL-26
        "%d/%B/%y",          # 19/APRIL/26
        "%d.%B.%y",          # 19.APRIL.26
        "%d %B %y",          # 19 APRIL 26

        "%B-%d-%y",          # APRIL-19-26
        "%B/%d-%y",          # APRIL/19-26
        "%B %d %y",          # APRIL 19 26
        "%B %d, %y",         # APRIL 19, 26

        "%Y%m%d",            # 20250711
        "%d%b%Y",            # 11Jul2025
        "%d%b%y",            # 11Jul25
        "%d%B%Y",            # 11July2025
        "%d%B%y",            # 11July25
    ]

    for fmt in date_formats:
        try:
            dt = datetime.strptime(date_str, fmt)
            if dt.year <= 99:
                dt = dt.replace(year=dt.year + 2000)
            return dt.strftime("%Y-%m-%d")
        except ValueError:
            continue

    # ── Excel serial number fallback (e.g. "45678") ───────────────────────
    try:
        serial = int(date_str)
        if 30000 <= serial <= 60000:          # sanity range: ~1982 – ~2064
            from datetime import timedelta
            excel_epoch = datetime(1899, 12, 30)
            dt = excel_epoch + timedelta(days=serial)
            return dt.strftime("%Y-%m-%d")
    except (ValueError, TypeError):
        pass

    return None


def normalize_attachment_text(text: str) -> str:
    if not text:
        return ""
    text = text.replace("\u2013", "-").replace("\u2014", "-")
    text = re.sub(r"\s*-\s*", "-", text)
    # DO NOT flatten newlines — they preserve table row/column structure
    # Fix broken word-wrap across lines (word- \n word → word - word)
    text = re.sub(r"([A-Za-z])-\n\s*([A-Za-z])", r"\1-\2", text)
    # Normalize dates like 2025-07-06
    text = re.sub(
        r"(\d{4})-(\d{1,2})-(\d{1,2})",
        lambda m: f"{m.group(1)}-{int(m.group(2)):02d}-{int(m.group(3)):02d}",
        text,
    )
    text = re.sub(r"[ \t]+", " ", text)       # collapse horizontal whitespace only
    text = re.sub(r"\n{3,}", "\n\n", text)    # max 2 consecutive blank lines
    return text.strip()


def extract(pattern: str, text: str) -> Optional[str]:
    match = re.search(pattern, text, re.IGNORECASE)
    return match.group(1).strip() if match else None

ITEM_BLOCK_REGEX = re.compile(
    r"""
    (
        .{0,200}?                                  # Description context
        (?:\b\d{2}K\b|\b\d{2}KT\b|\bAG\d{3}\b)     # Gold karat
        .{0,200}?                                  # More context
        (?:qty|quantity|\b\d+\b)                  # Quantity indicator
    )
    """,
    re.IGNORECASE | re.DOTALL | re.VERBOSE
)


def extract_po_items(text: str) -> list:
    """
    Extracts per-row items from structured attachment text.
    Works on the pipe-delimited row format produced by extract_text_from_attachment.
    Falls back to an empty list cleanly so the caller inserts a header-only record.
    """
    items = []

    for line in text.splitlines():
        line = line.strip()
        if not line or line.startswith("Sheet:"):
            continue

        karat_match = re.search(r"\b(\d{2}\s*K[TWY]?|\d{2}\s*KT|AG\d{3})\b", line, re.IGNORECASE)
        if not karat_match:
            continue

        qty_match   = re.search(r"(?:qty|quantity|pcs|ord#?)\s*[:\-]?\s*(\d+)|\b(\d+)\s*(?:pcs|pc|ea)\b", line, re.IGNORECASE)
        date_match  = re.search(r"\b(\d{4}-\d{2}-\d{2}|\d{1,2}[/-]\d{1,2}[/-]\d{2,4})\b", line)
        color_match = re.search(r"\b(yellow|white|rose)\b", line, re.IGNORECASE)
        desc_match  = re.search(r"([A-Za-z][A-Za-z0-9 ,\-]{9,})", line)

        qty_raw = qty_match.group(1) or qty_match.group(2) if qty_match else None

        items.append({
            "description":    desc_match.group(1).strip() if desc_match else None,
            "gold_karat":     karat_match.group(1).strip().upper(),
            "quantity":       int(qty_raw) if qty_raw else None,
            "delivery_date":  date_match.group(1) if date_match else None,
            "color":          color_match.group(1).capitalize() if color_match else None,
            "gold_lock":      None,
        })

    return items


async def extract_po_header(text: str):
    return await extract_po_fields(text)


IMAGE_EXTENSIONS = (".png", ".jpg", ".jpeg", ".bmp", ".tiff", ".webp")

# ---------------- TEXT EXTRACTION ---------------- #
def xlrd_cell_to_str(cell, workbook) -> str:
    import xlrd
    from datetime import datetime

    try:
        # Proper date handling
        if cell.ctype == xlrd.XL_CELL_DATE:
            dt = xlrd.xldate_as_datetime(cell.value, workbook.datemode)
            return dt.strftime("%Y-%m-%d")

        # Handle numbers that might actually be dates
        if cell.ctype == xlrd.XL_CELL_NUMBER:
            # Try converting to date if possible
            try:
                dt = xlrd.xldate_as_datetime(cell.value, workbook.datemode)
                return dt.strftime("%Y-%m-%d")
            except Exception:
                return str(cell.value)

        if cell.ctype == xlrd.XL_CELL_EMPTY:
            return ""

        return str(cell.value)

    except Exception:
        return str(cell.value)


async def extract_text_from_attachment(content_bytes, filename, content_type):
    ext = (filename or "").lower()
    ct = (content_type or "").lower()

    def parse_attachment():
        try:
            if ct.startswith("text/") or ext.endswith((".txt", ".md", ".csv", ".log")):
                return content_bytes.decode("utf-8", errors="ignore")

            elif ct == "application/pdf" or ext.endswith(".pdf"):
                import PyPDF2
                reader = PyPDF2.PdfReader(io.BytesIO(content_bytes))
                return "\n".join((p.extract_text() or "") for p in reader.pages)

            elif ct in ("application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                        "application/msword") or ext.endswith((".docx", ".doc")):
                import docx
                document = docx.Document(io.BytesIO(content_bytes))
                return "\n".join(p.text for p in document.paragraphs)

            elif ct in ("application/vnd.openxmlformats-officedocument.presentationml.presentation",
                        "application/vnd.ms-powerpoint") or ext.endswith((".pptx", ".ppt")):
                from pptx import Presentation
                prs = Presentation(io.BytesIO(content_bytes))
                return "\n".join(
                    shape.text
                    for slide in prs.slides
                    for shape in slide.shapes
                    if hasattr(shape, "text")
                )

            elif ct in (
                "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                "application/vnd.ms-excel"
            ) or ext.endswith((".xlsx", ".xls")):
                all_text = []
                if ext.endswith(".xlsx"):
                    from openpyxl import load_workbook
                    from datetime import datetime

                    wb = load_workbook(io.BytesIO(content_bytes), data_only=True)

                    for sheet in wb.worksheets:   
                        all_text.append(f"Sheet: {sheet.title}")

                        for row in sheet.iter_rows():
                            values = []
                            for cell in row:
                                val = cell.value

                                # Proper date handling
                                if isinstance(val, datetime):
                                    val = val.strftime("%Y-%m-%d")

                                values.append(str(val) if val is not None else "")

                            row_text = " | ".join(v for v in values if v.strip())
                            if row_text:
                                all_text.append(row_text)
                elif ext.endswith(".xls"):
                    import xlrd
                    wb = xlrd.open_workbook(file_contents=content_bytes)
                    for sheet in wb.sheets():
                        all_text.append(f"Sheet: {sheet.name}")
                        for row_idx in range(sheet.nrows):
                            row_vals = [
                                xlrd_cell_to_str(sheet.cell(row_idx, col_idx), wb)
                                for col_idx in range(sheet.ncols)
                            ]
                            row_text = " | ".join(v for v in row_vals if v.strip())
                            if row_text.strip():
                                all_text.append(row_text)
                return "\n".join(all_text)

            elif any(ext.endswith(e) for e in IMAGE_EXTENSIONS):
                try:
                    return extract_text_from_image_bytes(content_bytes)
                except Exception as e:
                    logger.error(f"OCR failed for {filename}: {e}")
                    return ""

        except Exception as e:
            logger.error(f"parse_attachment failed for {filename}: {e}")
            return None

    return await asyncio.to_thread(parse_attachment)


# -------make subject line detection of RE/FW(handles multiple prefixes, spaces, case variations)---------
PREFIX_PATTERN = re.compile(r"^\s*(re|fw|fwd)\s*:", re.IGNORECASE)

def is_reply_or_forward(subject: str) -> bool:
    """ Detects RE / FW / FWD prefixes in a safe, case-insensitive way. """
    if not subject:
        return False

    try:
        original = subject.strip()
        temp = original

        # Remove chained prefixes like RE: FW: RE:
        while PREFIX_PATTERN.match(temp):
            temp = PREFIX_PATTERN.sub("", temp, count=1).strip()

        # If changed → it had RE/FW prefix
        return temp != original

    except Exception:
        # Never break main flow
        return False
    
    
# ------------------- Main Function to Fetch and Save Emails + Attachments ------------------- #
async def fetch_and_save_mails_by_folders(
    access_token: str,
    folder_names: list[str],
    user_id: int,
    from_date: str,
    to_date: str,
    mails_repo: "MailsRepository"
) -> List[Dict[str, Any]]:

    extracted_po_ids: list[int] = []
    headers = {"Authorization": f"Bearer {access_token}"}
    results: List[Dict[str, Any]] = []

    from_date_iso = f"{from_date}T00:00:00Z"
    to_date_iso = f"{to_date}T23:59:59Z"

    timeout = httpx.Timeout(connect=10.0, read=300.0, write=60.0, pool=30.0)

    async with httpx.AsyncClient(timeout=timeout) as client:

        # ---------------- FETCH FOLDERS ----------------
        try:
            folder_resp = await client.get(
                f"{GRAPH_API}/me/mailFolders?$top=200&$expand=childFolders",
                headers=headers
            )
            folder_resp.raise_for_status()
            logger.info(f"folder response:{folder_resp}")
        except Exception as e:
            logger.error(f"Failed to fetch folders:{e}")
            return []

        folders = folder_resp.json().get("value", [])
        wanted = {f.lower() for f in folder_names}

        # ---------------- PROCESS FOLDERS ----------------
        for folder in folders:
            folder_id = folder.get("id")
            folder_name = folder.get("displayName", "")

            if not folder_id or folder_name.lower() not in wanted:
                continue

            url = (
                f"{GRAPH_API}/me/mailFolders/{folder_id}/messages"
                f"?$filter=receivedDateTime ge {from_date_iso} and receivedDateTime le {to_date_iso}"
                f"&$select=id,subject,body,from,toRecipients,ccRecipients,bccRecipients,"
                f"hasAttachments,receivedDateTime,bodyPreview"
            )

            messages = []
            next_url = url

            keywords = await mails_repo.fetch_keywords()
            logger.info(f"keywords are:{keywords}")

            # ---------------- PAGINATION ----------------
            while next_url:
                try:
                    resp = await client.get(next_url, headers=headers)
                    resp.raise_for_status()
                except Exception as e:
                    logger.error("Failed to fetch messages from folder %s: %s", folder_name, e)
                    break

                data = resp.json()
                messages.extend(data.get("value", []))
                next_url = data.get("@odata.nextLink")


            # ---------------- PROCESS EACH MESSAGE ----------------
            for msg in messages:
                graph_mail_id = msg.get("id")
                if not graph_mail_id:
                    continue

                if msg.get('@odata.type') in (
                    '#microsoft.graph.eventMessage',
                    '#microsoft.graph.eventMessageRequest'
                ):
                    continue

                if await mails_repo.mail_exists(graph_mail_id, user_id):
                    continue

                # ------------- CHECK DOMAIN FIRST -------------
                from_email = (msg.get('from') or {}).get('emailAddress', {}).get('address')

                sender_domain = None

                if from_email and "@" in from_email:
                    sender_domain = from_email.split("@")[-1].lower().strip()

                    if sender_domain in BLOCKED_DOMAINS:
                        logger.info(f"Skipping mail from blocked domain: {from_email}")
                        continue

                #-----------process body-----------
                subject = msg.get('subject', '')

                #-----------Skip reply/forward mails----------- 
                if is_reply_or_forward(subject):
                    logger.info(f"Skipping reply/forward mail: {subject}")
                    continue

                body_obj = msg.get('body') or {}
                body_content = body_obj.get('content') or msg.get('bodyPreview', '')

                body_plain = strip_html_to_text(body_content)
                body_clean = clean_email_body(body_plain)

                # ---------------- LLM KEYWORD MATCH ----------------
                matched_keywords, match_source = await detect_keywords(
                    body_clean,
                    keywords
                )
                logger.info(f"Matched Keywords are:{matched_keywords}")

                if not matched_keywords:
                    continue  #skip mail

                from_email = (msg.get('from') or {}).get('emailAddress', {}).get('address')
                to_emails = collect_addresses_from_message(msg, 'toRecipients')
                cc = collect_addresses_from_message(msg, 'ccRecipients')
                bcc = collect_addresses_from_message(msg, 'bccRecipients')
                merged_cc = ",".join(filter(None, [cc, bcc])) or None
                has_attachments = bool(msg.get('hasAttachments'))
                date_only = iso_to_date(msg.get('receivedDateTime'))

                # ---------------- INSERT MAILS ----------------
                try:
                    mail_id = await mails_repo.insert_mail_detail(
                    user_id=user_id,
                    subject=subject,
                    body=body_clean,
                    date_time=date_only,
                    mail_from=from_email,
                    mail_to=to_emails,
                    mail_cc=merged_cc,
                    created_by=user_id,
                    updated_by=user_id,
                    is_active=1,
                    graph_mail_id=graph_mail_id,
                    folder_name=folder_name
                )
                    logger.info(f"Mail id:{mail_id}")
                except Exception as e:
                    logger.error(f"DB insert failed: {str(e)}")
                    continue

                saved_attachments = []
                attachment_texts = []

                # ---------------- ATTACHMENTS ----------------
                if has_attachments:
                    try:
                        att_resp = await client.get(
                            f"{GRAPH_API}/me/messages/{graph_mail_id}/attachments",
                            headers=headers
                        )
                        logger.info(f"attchment response:{att_resp}")
                        att_resp.raise_for_status()
                        att_list = att_resp.json().get("value", [])
                    except Exception as e:
                        logger.error("Failed fetching attachments for %s: %s", graph_mail_id, e)
                        att_list = []

                    for att in att_list:
                        filename = att.get("name")
                        content_type = (att.get("contentType") or "").lower()
                        content_bytes = base64.b64decode(att.get("contentBytes") or "")
                        if not content_bytes:
                            continue

                        file_hash = compute_file_hash(content_bytes)
                        if await mails_repo.attachment_exists(file_hash,user_id):
                            continue

                        # Save File
                        safe_filename = re.sub(r'[\\/*?:"<>|&]', "_", filename)
                        os.makedirs("attachments", exist_ok=True)
                        file_path = os.path.join("attachments", safe_filename)
                        with open(file_path, "wb") as f:
                            f.write(content_bytes)

                        attachment_text = await extract_text_from_attachment(content_bytes, filename, content_type)
                        if not attachment_text:
                            continue

                        attach_keywords, match_type = await detect_keywords(attachment_text, keywords)

                        if not attach_keywords:
                            logger.info(f"Skipping attachment '{filename}' — no keyword match.")
                            continue

                        # Only add to attachment_texts AFTER keyword check passes
                        attachment_texts.append(attachment_text)

                        try:
                            response = await mails_repo.insert_attachment(
                                mail_dtl_id=mail_id,
                                user_id=user_id,
                                attach_name=filename,
                                attach_type=content_type,
                                attach_path=file_path,
                                created_by=user_id,
                                updated_by=user_id,
                                is_active=1,
                                file_hash=file_hash,
                            )
                            logger.info(f"response is:{response}")
                            saved_attachments.append(filename)
                        except Exception as e:
                            logger.error("Attachment insert failed (%s): %s", filename, e)

                
                # ---------------- Insert PO data from email body ----------------
                po_data_body = await extract_po_fields(body_clean)
                logger.info(f"PO Body Data:{po_data_body}")
                if po_data_body.get("po_number") or po_data_body.get("customer_name"):
                    po_det_id = await mails_repo.insert_po_details(
                        mail_dtl_id=mail_id,
                        user_id=user_id,
                        po_number=po_data_body.get("po_number"),
                        customer_name=po_data_body.get("customer_name"),
                        vendor_number=po_data_body.get("vendor_number"),
                        po_date=normalize_po_date_ddmmyyyy(po_data_body.get("po_date")),
                        delivery_date=normalize_po_date_ddmmyyyy(po_data_body.get("delivery_date")),
                        cancel_date=normalize_po_date_ddmmyyyy(po_data_body.get("cancel_date")),
                        gold_karat=po_data_body.get("gold_karat"),
                        ec_style_number=po_data_body.get("ec_style_number"),
                        customer_style_number=po_data_body.get("customer_style_number"),
                        color=po_data_body.get("color"),
                        quantity=po_data_body.get("quantity"),
                        description=po_data_body.get("description"),
                        mail_folder=folder_name,
                        created_by=user_id,
                        gold_lock=po_data_body.get("gold_lock"),
                        domain_name=sender_domain,
                        source="email"
                    )
                    extracted_po_ids.append(po_det_id)
                    logger.info(f"Inserted Mail PO Extracted IDs:{extracted_po_ids}")
                # ----------------Insert PO data from attachments ----------------
                for att_text in attachment_texts:
                    normalized_text = normalize_attachment_text(att_text)
                    # ---------------- PO HEADER FROM ATTACHMENT ----------------
                    header = await extract_po_header(normalized_text)

                    if not any(header.values()):
                        continue

                    # ---------------- PO ITEMS FROM ATTACHMENT ----------------
                    # extract_po_items(normalized_text)
                    items = None
                    logger.info(f"Items:{items}")

                    # Fallback: if no items found, insert header-only
                    if not items:
                        po_det_id = await mails_repo.insert_po_details(
                            mail_dtl_id=mail_id,
                            user_id=user_id,
                            po_number=header.get("po_number"),
                            customer_name=header.get("customer_name"),
                            vendor_number=header.get("vendor_number"),
                            po_date=normalize_po_date_ddmmyyyy(header.get("po_date")),
                            delivery_date=normalize_po_date_ddmmyyyy(header.get("delivery_date")),
                            cancel_date=normalize_po_date_ddmmyyyy(header.get("cancel_date")),
                            gold_karat=header.get("gold_karat"),
                            ec_style_number=header.get("ec_style_number"),
                            customer_style_number=header.get("customer_style_number"),
                            color=header.get("color"),
                            quantity=header.get("quantity"),
                            description=header.get("description"),
                            mail_folder=folder_name,
                            created_by=user_id,
                            gold_lock=header.get("gold_lock"),
                            domain_name=sender_domain,
                            source="email"
                        )
                        extracted_po_ids.append(po_det_id)
                        logger.info(f"Inserted Attchment PO Extracted IDs:{extracted_po_ids}")
                    else:
                        # MULTIPLE ROW INSERTS
                        for item in items:
                            po_det_id = await mails_repo.insert_po_details(
                                mail_dtl_id=mail_id,
                                user_id=user_id,
                                po_number=header.get("po_number"),
                                customer_name=header.get("customer_name"),
                                vendor_number=header.get("vendor_number"),
                                po_date=normalize_po_date_ddmmyyyy(header.get("po_date")),
                                delivery_date=normalize_po_date_ddmmyyyy(item.get("delivery_date")),
                                cancel_date=normalize_po_date_ddmmyyyy(header.get("cancel_date")),
                                gold_karat=item.get("gold_karat"),
                                ec_style_number=header.get("ec_style_number"),
                                customer_style_number=header.get("customer_style_number"),
                                color=header.get("color"),
                                quantity=item.get("quantity"),
                                description=item.get("description"),
                                mail_folder=folder_name,
                                created_by=user_id,
                                gold_lock=header.get("gold_lock"),
                                domain_name=sender_domain,
                                source="email"
                            )
                            extracted_po_ids.append(po_det_id)
                            logger.info(f"Inserted Attachment PO Extracted IDs:{extracted_po_ids}")
                # ---------------- COLLECT RESULT ----------------
                results.append({
                    "mail_dtl_id": mail_id,
                    "subject": subject,
                    "from": from_email,
                    "to": to_emails,
                    "cc": merged_cc,
                    "has_attachments": has_attachments,
                    "attachments": saved_attachments,
                    "folder": folder_name,
                })

    return {
    "results": results,
    "extracted_po_ids": extracted_po_ids
    }
# ------------------Email + Attachment Fetching + LLM logic end ------------------ #


# ============================================================
# data comparison logic start
# ============================================================
FIELDS_TO_COMPARE = [
    "customer_name",
    "vendor_number",
    "po_date",
    "po_number",
    "delivery_date",
    "cancel_date",
    "gold_lock",
    "ec_style_number",
    "customer_style_number",
    "gold_karat",
    "color",
    "quantity",
    "description",
]


def make_json_safe(obj):
    if isinstance(obj, (date, datetime)):
        return obj.isoformat()
    if isinstance(obj, Decimal):
        return float(obj)
    if isinstance(obj, bytes):
        try:
            return obj.decode("utf-8")
        except UnicodeDecodeError:
            return base64.b64encode(obj).decode("utf-8")
    return obj

# ============================================================
# LLM comparison function with robust JSON parsing and error handling
# ===========================================================
async def llm_batch_compare(pairs_for_llm):

    prompt = f"""
You are an expert PO field comparison engine.

Determine if the scanned value and system value represent the SAME meaning.

Rules:
- Treat spelling mistakes as SAME
- Treat abbreviations as SAME
- Treat casing differences as SAME
- Treat punctuation differences as SAME
- Treat word order differences as SAME

Examples considered SAME:
"yellow color" vs "color yellow"
"intl corp" vs "international corporation"
"ring size 7" vs "size 7 ring"

Only report mismatches if the meaning is clearly different.

Return ONLY JSON:
[
  {{
    "po_det_id": number,
    "system_po_id": number,
    "field": string,
    "scanned_value": string,
    "system_value": string
  }}
]

PO field data:
{json.dumps(pairs_for_llm, indent=2)}
"""

    resp = openai_client.chat.completions.create(
        model="gpt-4.1-mini",
        temperature=0,
        messages=[{"role": "user", "content": prompt}]
    )

    raw = resp.choices[0].message.content.strip()
    raw = re.sub(r"^```json|```$", "", raw, flags=re.IGNORECASE).strip()

    try:
        result = json.loads(raw)
        return result if isinstance(result, list) else []
    except Exception as e:
        print(f"LLM parsing error: {e}")
        return []
    

def chunk(data, size):
    for i in range(0, len(data), size):
        yield data[i:i + size]

def normalize_po(po):
    if not po:
        return ""
    return re.sub(r'[^A-Za-z0-9]', '', po).lower()


def normalize_value(v):
    if v is None:
        return ""

    v = str(v).lower().strip()
    # remove punctuation
    v = re.sub(r'[^a-z0-9\s]', '', v)
    # collapse spaces
    v = re.sub(r'\s+', ' ', v)

    return v

# ============================================================
# PO Recomparison start
# ============================================================
import hashlib

# hash on po_number + domain to stay consistent with unique key
def make_stable_system_po_id(po: dict) -> int:
    raw = "|".join([
        normalize_po(str(po.get("po_number") or "")),
        normalize_domain(str(po.get("domain_name") or "")), 
        str(po.get("order_date") or ""),
    ])
    full_hash = hashlib.sha256(raw.encode()).digest()
    return int.from_bytes(full_hash[:4], byteorder="big")


def normalize_domain(domain):
    if not domain:
        return ""
    return domain.lower().strip()


# ============================================================
# both po_number+domain must be present and match for any reconciliation to occur
# ============================================================
async def reconcile_all_pos(user_id, mails_repo, system_pos, batch_size=500):

    stats = {
        "mismatch_checked": 0,
        "mismatch_resolved": 0,
        "missing_checked": 0,
        "missing_resolved": 0,
        "errors": []
    }

    try:
        # ── Build system lookup with SAME key as compare function: (po_number, domain)
        system_po_map = defaultdict(list)
        for po in system_pos:
            po_number   = normalize_po(po.get("po_number"))
            domain_name = normalize_domain(po.get("domain_name"))
            if po_number and domain_name:                      # both must exist
                key = (po_number, domain_name)
                system_po_map[key].append(po)

        try:
            mismatch_stats = await reconcile_mismatches(user_id, mails_repo, system_po_map, batch_size)
            stats.update(mismatch_stats)
        except Exception as e:
            logger.error(f"[RECONCILE] reconcile_mismatches failed: {e}", exc_info=True)
            stats["errors"].append(f"mismatch: {str(e)}")

        try:
            missing_stats = await reconcile_missing(user_id, mails_repo, system_po_map, batch_size)
            stats.update(missing_stats)
        except Exception as e:
            logger.error(f"[RECONCILE] reconcile_missing failed: {e}", exc_info=True)
            stats["errors"].append(f"missing: {str(e)}")

    except Exception as e:
        logger.error(f"[RECONCILE] reconcile_all_pos outer exception: {e}", exc_info=True)
        stats["errors"].append(str(e))

    return stats

# ============================================================
# SHARED HELPER 
# ============================================================
async def resolve_all_differences(scanned: dict, system: dict) -> list:
    """
    Compares all 13 fields between scanned and system PO.
    Uses LLM for value-vs-value differences (same as compare function).
    Returns list of real mismatches. Empty list = all fields match.
    """
    pairs_for_llm   = []
    direct_mismatch = []

    for field in FIELDS_TO_COMPARE:
        scanned_val = scanned.get(field)
        system_val  = system.get(field)

        scanned_empty = scanned_val in (None, "")
        system_empty  = system_val  in (None, "")

        # Both empty → no issue
        if scanned_empty and system_empty:
            continue

        # System has value, scanned empty → direct mismatch, no LLM needed
        if not system_empty and scanned_empty:
            direct_mismatch.append({
                "po_det_id":     scanned["po_det_id"],
                "system_po_id":  system["system_po_id"],
                "field":         field,
                "scanned_value": "",
                "system_value":  str(system_val)
            })
            continue

        # Both have values but differ → send to LLM
        if normalize_value(scanned_val) != normalize_value(system_val):
            pairs_for_llm.append({
                "po_det_id":     scanned["po_det_id"],
                "system_po_id":  system["system_po_id"],
                "field":         field,
                "scanned_value": str(scanned_val),
                "system_value":  str(system_val)
            })

    # LLM decides on value-vs-value differences
    llm_mismatches = []
    if pairs_for_llm:
        try:
            llm_mismatches = await llm_batch_compare(pairs_for_llm)
        except Exception as e:
            logger.error(f"LLM compare failed: {e} — treating all as mismatches")
            llm_mismatches = pairs_for_llm  # safe fallback

    return llm_mismatches + direct_mismatch

# ===========================================================
# reconcile_mismatches — for active mismatches, check if they can be resolved based on current system data
# =========================================================
async def reconcile_mismatches(user_id, mails_repo, system_po_map, batch_size):
    stats = {"mismatch_checked": 0, "mismatch_resolved": 0}

    active = await mails_repo.get_all_active_mismatches(user_id)
    if not active:
        return stats

    stats["mismatch_checked"] = len(active)

    grouped = defaultdict(list)
    for row in active:
        grouped[row["po_det_id"]].append(row)

    po_det_ids   = list(grouped.keys())
    scanned_list = await mails_repo.get_po_details_by_ids(po_det_ids)
    scanned_map  = {p["po_det_id"]: p for p in scanned_list}

    for batch_keys in chunk(po_det_ids, batch_size):
        for po_det_id in batch_keys:

            scanned = scanned_map.get(po_det_id)
            if not scanned or not is_valid_for_matching(scanned):
                continue

            key = (normalize_po(scanned.get("po_number")), normalize_domain(scanned.get("domain_name")))
            candidates = system_po_map.get(key)
            if not candidates:
                continue

            system = find_exact_system_match(scanned, candidates)
            if not system:
                continue

            # Update customer_name from system if different
            system_customer_name = system.get("customer_name")
            if system_customer_name and normalize_value(scanned.get("customer_name")) != normalize_value(system_customer_name):
                try:
                    await mails_repo.update_po_customer_name(po_det_id=scanned["po_det_id"], customer_name=system_customer_name)
                    scanned["customer_name"] = system_customer_name
                except Exception as e:
                    logger.error(f"Failed to update customer_name for po_det_id={scanned.get('po_det_id')}: {e}")

            try:
                # Same logic as compare — check all 13 fields with LLM
                all_mismatches = await resolve_all_differences(scanned, system)

                if all_mismatches:
                    # still has differences → stay in mismatch, do nothing
                    logger.info(f"po_det_id={po_det_id} → still mismatched ({len(all_mismatches)} field(s))")
                    continue

                # all fields now match → deactivate mismatch + insert matched
                mismatch_ids = [r["po_mismatch_id"] for r in grouped[po_det_id]]
                await mails_repo.deactivate_mismatches(user_id, mismatch_ids)

                exists = await mails_repo.matched_po_exists(user_id, po_det_id, system["system_po_id"])
                if not exists:
                    await mails_repo.insert_matched_po(
                        po_det_id=po_det_id,
                        system_po_id=system["system_po_id"],
                        mail_dtl_id=scanned.get("mail_dtl_id"),
                        user_id=user_id,
                        po_number=scanned.get("po_number"),
                        po_date=scanned.get("po_date"),
                        vendor_number=scanned.get("vendor_number"),
                        customer_name=scanned.get("customer_name"),
                        created_by="reconciliation"
                    )

                stats["mismatch_resolved"] += 1
                logger.info(f"po_det_id={po_det_id} → mismatch resolved → matched")

            except Exception as e:
                logger.error(f"Mismatch resolve error: {e}")

    return stats


# ============================================================
# reconcile_missing 
# ============================================================
async def reconcile_missing(user_id, mails_repo, system_po_map, batch_size):
    stats = {"missing_checked": 0, "missing_resolved": 0, "mismatch_created": 0}

    active = await mails_repo.get_all_active_missing(user_id)
    if not active:
        return stats

    stats["missing_checked"] = len(active)

    po_det_ids   = list({r["po_det_id"] for r in active if r.get("po_det_id")})
    scanned_list = await mails_repo.get_po_details_by_ids(po_det_ids)
    scanned_map  = {p["po_det_id"]: p for p in scanned_list}

    for batch in chunk(active, batch_size):
        for row in batch:

            po_det_id  = row["po_det_id"]
            missing_id = row["po_missing_id"]

            scanned = scanned_map.get(po_det_id)
            if not scanned or not is_valid_for_matching(scanned):
                continue

            key = (normalize_po(scanned.get("po_number")), normalize_domain(scanned.get("domain_name")))
            candidates = system_po_map.get(key)

            # still not found in system → stay missing
            if not candidates:
                continue

            system = find_exact_system_match(scanned, candidates)
            if not system:
                continue

            # Update customer_name from system if different
            system_customer_name = system.get("customer_name")
            if system_customer_name and normalize_value(scanned.get("customer_name")) != normalize_value(system_customer_name):
                try:
                    await mails_repo.update_po_customer_name(po_det_id=scanned["po_det_id"], customer_name=system_customer_name)
                    scanned["customer_name"] = system_customer_name
                except Exception as e:
                    logger.error(f"Failed to update customer_name for po_det_id={scanned.get('po_det_id')}: {e}")

            try:
                # Same logic as compare — check all 13 fields with LLM
                all_mismatches = await resolve_all_differences(scanned, system)

                if not all_mismatches:
                    # all fields match → deactivate missing + insert matched
                    await mails_repo.deactivate_missing_pos(user_id, missing_id)

                    exists = await mails_repo.matched_po_exists(user_id, po_det_id, system["system_po_id"])
                    if not exists:
                        await mails_repo.insert_matched_po(
                            po_det_id=po_det_id,
                            system_po_id=system["system_po_id"],
                            mail_dtl_id=scanned.get("mail_dtl_id"),
                            user_id=user_id,
                            po_number=scanned.get("po_number"),
                            po_date=scanned.get("po_date"),
                            vendor_number=scanned.get("vendor_number"),
                            customer_name=scanned.get("customer_name"),
                            created_by="reconciliation"
                        )

                    stats["missing_resolved"] += 1
                    logger.info(f"po_det_id={po_det_id} → missing resolved → matched")

                else:
                    # PO found but fields differ → deactivate missing + insert mismatches
                    # Always deactivate missing — PO was found in system, it is no longer missing
                    for mm in all_mismatches:
                        scanned_value = str(mm["scanned_value"]) if mm["scanned_value"] else ""
                        system_value  = str(mm["system_value"])  if mm["system_value"]  else ""

                        exists = await mails_repo.mismatch_exists(
                            user_id=user_id,
                            po_det_id=po_det_id,
                            system_po_id=system["system_po_id"],
                            mismatch_attribute=mm["field"],
                            scanned_value=scanned_value,
                            system_value=system_value
                        )
                        if not exists:
                            await mails_repo.insert_mismatch(
                                po_det_id=po_det_id,
                                user_id=user_id,
                                system_po_id=system["system_po_id"],
                                field=mm["field"],
                                scanned_value=scanned_value,
                                system_value=system_value,
                                comment=f"{mm['field']} mismatch"
                            )

                    # Always deactivate missing regardless of whether mismatches were new or existing
                    await mails_repo.deactivate_missing_pos(user_id, missing_id)
                    stats["mismatch_created"] += 1
                    logger.info(f"po_det_id={po_det_id} → missing → mismatch ({len(all_mismatches)} field(s))")

            except Exception as e:
                logger.error(f"Reconcile missing error: {e}")

    return stats

# ============================================================
# fetch oldest date to optimize the system PO fetch for comparison/recomparison
# ============================================================
async def fetch_system_pos_with_oldest_date(mails_repo, app):
    oldest_date = await mails_repo.get_oldest_report_date()

    if oldest_date:
        system_pos = await MSSQLRepo.get_po_list(app, oldest_date)
    else:
        system_pos = await MSSQLRepo.get_po_list_without_oldest_date(app)

    for po in system_pos:
        po["system_po_id"] = make_stable_system_po_id(po)

    system_pos = [
        {k: make_json_safe(v) for k, v in po.items()}
        for po in system_pos
    ]

    return system_pos


# ============================================================
# CORE: Strong match = po_number + domain_name (unique key)
# If either is missing → go to missing directly
# ============================================================
def is_valid_for_matching(scanned: dict) -> bool:
    """Both po_number AND domain_name must exist to attempt matching."""
    return bool(normalize_po(scanned.get("po_number"))) and \
           bool(normalize_domain(scanned.get("domain_name")))


def find_exact_system_match(scanned: dict, candidates: list) -> dict | None:
    """
    Strong exact match on po_number + domain_name (unique key).
    Returns the matched system PO or None.
    No fuzzy scoring — this is a hard gate.
    """
    scanned_po     = normalize_po(scanned.get("po_number"))
    scanned_domain = normalize_domain(scanned.get("domain_name"))
    mail_date      = scanned.get("date_time")

    # Parse mail date once
    parsed_mail_date = None
    if mail_date:
        if isinstance(mail_date, (date, datetime)):
            parsed_mail_date = mail_date.date() if isinstance(mail_date, datetime) else mail_date
        else:
            try:
                parsed_mail_date = datetime.strptime(str(mail_date)[:10], "%Y-%m-%d").date()
            except Exception:
                parsed_mail_date = None

    for system in candidates:
        system_po     = normalize_po(system.get("po_number"))
        system_domain = normalize_domain(system.get("domain_name"))

        # ── po_number + domain must match
        if scanned_po != system_po or scanned_domain != system_domain:
            continue
        
        # ── EMR order_date must be >= mail received date ──
        if parsed_mail_date:
            order_date_raw = system.get("order_date")
            parsed_order_date = None

            if order_date_raw:
                if isinstance(order_date_raw, (date, datetime)):
                    parsed_order_date = order_date_raw.date() if isinstance(order_date_raw, datetime) else order_date_raw
                else:
                    try:
                        parsed_order_date = datetime.strptime(str(order_date_raw)[:10], "%Y-%m-%d").date()
                    except Exception:
                        parsed_order_date = None

            if parsed_order_date and parsed_order_date < parsed_mail_date:
                logger.info(
                    f"Skipping system PO (order_date={parsed_order_date} < mail_date={parsed_mail_date}) "
                    f"for po_number={scanned.get('po_number')}"
                )
                continue # reject — EMR order is older than the mail
        
        return system  # passed all checks

    return None  # no exact match and EMR order_date >= mail received date  → treat as missing

# ============================================================
# compare_scanned_and_system_pos 
# ============================================================
async def compare_scanned_and_system_pos(
    request=None,
    app=None,
    user_id: int = None,
    po_det_ids: list[int] = None,
    mails_repo=None,
    system_pos=None
):
    try:
        resolved_app = app or (request.app if request is not None else None)

        if resolved_app is None:
            logger.error("No app context provided to compare_scanned_and_system_pos")
            return {
                "status": "error",
                "message": "Failed to generate PO report",
                "error": "No app context available"
            }

        scanned_pos = await mails_repo.get_po_details_by_ids(po_det_ids)

        if not scanned_pos:
            return {"status": "success", "message": "No scanned POs found for comparison"}

        scanned_pos = [
            {k: make_json_safe(v) for k, v in po.items()}
            for po in scanned_pos
        ]

        # Build system PO lookup: key = (po_number, domain_name) — unique key
        system_po_map = defaultdict(list)
        for po in system_pos:
            po_number   = normalize_po(po.get("po_number"))
            domain_name = normalize_domain(po.get("domain_name"))
            if po_number and domain_name:
                key = (po_number, domain_name)
                system_po_map[key].append(po)

        matched_pairs = []
        missing_pos   = []

        for scanned in scanned_pos:

            # if po_number OR domain missing → missing directly, no matching attempt
            if not is_valid_for_matching(scanned):
                logger.info(
                    f"po_det_id={scanned.get('po_det_id')} → missing "
                    f"(po_number={scanned.get('po_number')}, domain={scanned.get('domain_name')})"
                )
                missing_pos.append(scanned)
                continue

            key = (normalize_po(scanned.get("po_number")), normalize_domain(scanned.get("domain_name")))
            candidates = system_po_map.get(key)

            # No candidates found for this unique key → missing
            if not candidates:
                logger.info(f"po_det_id={scanned.get('po_det_id')} → missing (no system PO for key={key})")
                missing_pos.append(scanned)
                continue

            # Exact match on unique key (po_number + domain)
            system = find_exact_system_match(scanned, candidates)

            if not system:
                logger.info(f"po_det_id={scanned.get('po_det_id')} → missing (exact match failed for key={key})")
                missing_pos.append(scanned)
                continue

            # Update customer_name if system has a cleaner version
            system_customer_name = system.get("customer_name")
            if (
                system_customer_name
                and normalize_value(scanned.get("customer_name")) != normalize_value(system_customer_name)
            ):
                try:
                    await mails_repo.update_po_customer_name(
                        po_det_id=scanned["po_det_id"],
                        customer_name=system_customer_name
                    )
                    scanned["customer_name"] = system_customer_name
                except Exception as e:
                    logger.error(f"Failed to update customer_name for po_det_id={scanned['po_det_id']}: {e}")

            matched_pairs.append({
                "po_det_id":   scanned["po_det_id"],
                "system_po_id": system["system_po_id"],
                "scanned":     {f: scanned.get(f) for f in FIELDS_TO_COMPARE},
                "system":      {f: system.get(f)  for f in FIELDS_TO_COMPARE},
                "raw_scanned": scanned
            })

        pairs_for_llm = []

        for pair in matched_pairs:
            for field in FIELDS_TO_COMPARE:
                scanned_val = pair["scanned"].get(field)
                system_val  = pair["system"].get(field)

                if scanned_val in (None, "") or system_val in (None, ""):
                    continue

                if normalize_value(scanned_val) == normalize_value(system_val):
                    continue

                pairs_for_llm.append({
                    "po_det_id":     pair["po_det_id"],
                    "system_po_id":  pair["system_po_id"],
                    "field":         field,
                    "scanned_value": str(scanned_val),
                    "system_value":  str(system_val)
                })

        mismatches = []
        if pairs_for_llm:
            mismatches = await llm_batch_compare(pairs_for_llm)

        mismatch_pairs = set()

        for mm in mismatches:
            scanned_value = "" if mm["scanned_value"] is None else str(mm["scanned_value"])
            system_value  = "" if mm["system_value"]  is None else str(mm["system_value"])

            exists = await mails_repo.mismatch_exists(
                user_id=user_id,
                po_det_id=mm["po_det_id"],
                system_po_id=mm["system_po_id"],
                mismatch_attribute=mm["field"],
                scanned_value=scanned_value,
                system_value=system_value
            )

            if not exists:
                await mails_repo.insert_mismatch(
                    po_det_id=mm["po_det_id"],
                    user_id=user_id,
                    system_po_id=mm["system_po_id"],
                    field=mm["field"],
                    system_value=system_value,
                    scanned_value=scanned_value,
                    comment=f"{mm['field']} mismatch"
                )

            mismatch_pairs.add((mm["po_det_id"], mm["system_po_id"]))

        for pair in matched_pairs:
            for field in FIELDS_TO_COMPARE:
                scanned_val = pair["scanned"].get(field)
                system_val  = pair["system"].get(field)

                if system_val not in (None, "") and scanned_val in (None, ""):
                    exists = await mails_repo.mismatch_exists(
                        user_id=user_id,
                        po_det_id=pair["po_det_id"],
                        system_po_id=pair["system_po_id"],
                        mismatch_attribute=field,
                        scanned_value="",
                        system_value=str(system_val)
                    )

                    if not exists:
                        await mails_repo.insert_mismatch(
                            po_det_id=pair["po_det_id"],
                            user_id=user_id,
                            system_po_id=pair["system_po_id"],
                            field=field,
                            system_value=str(system_val),
                            scanned_value="",
                            comment=f"{field} missing in scanned data"
                        )

                    mismatch_pairs.add((pair["po_det_id"], pair["system_po_id"]))

        for pair in matched_pairs:
            key = (pair["po_det_id"], pair["system_po_id"])
            if key in mismatch_pairs:
                continue

            scanned = pair["raw_scanned"]
            await mails_repo.insert_matched_po(
                po_det_id=pair["po_det_id"],
                system_po_id=pair["system_po_id"],
                mail_dtl_id=scanned.get("mail_dtl_id"),
                user_id=user_id,
                po_number=scanned.get("po_number"),
                po_date=scanned.get("po_date"),
                vendor_number=scanned.get("vendor_number"),
                customer_name=scanned.get("customer_name"),
                created_by="system"
            )

        for po in missing_pos:
            exists = await mails_repo.po_missing_exists(
                user_id=user_id,
                po_det_id=po["po_det_id"],
                system_po_id=None,
                mismatch_attribute="po_missing",
                scanned_value=po.get("po_number"),
                system_value=""
            )

            if not exists:
                await mails_repo.insert_po_missing(
                    po_det_id=po["po_det_id"],
                    user_id=user_id,
                    system_po_id=None,
                    attribute="po_missing",
                    system_value="",
                    scanned_value=po.get("po_number"),
                    comment="PO not found in system"
                )

        return {"status": "success", "message": "PO comparison completed successfully"}

    except Exception as e:
        logger.exception(f"Error in compare_scanned_and_system_pos | user_id={user_id}")
        return {"status": "error", "message": "Failed to generate PO report", "error": str(e)}
