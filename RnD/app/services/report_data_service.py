import re
import PyPDF2
import io
import csv
import os
import logging
from typing import List, Union, Dict, Tuple, Any, Optional
from app.db.repositories.report_data_repo import ReportDataRepo
from datetime import date, datetime
from collections import Counter
from fastapi import HTTPException
from docx import Document
from pptx import Presentation
from openpyxl import load_workbook
import xlrd

# Set up logging
logger = logging.getLogger(__name__)


class ReportDataService:
    def __init__(self, repository: ReportDataRepo):
        self.repository = repository

    async def process_mail_effort(
        self,
        user_id: int,
        org_id: int,
        mail_dtl_id: Optional[int] = None,
        report_id: Optional[int] = None,
        efforts: Optional[float] = None,
        keyword_efforts: Optional[int] = None,
    ) -> Dict[str, Any]:
        try:
            # --- Update existing report ---
            if report_id:
                await self.repository.update_report_data(
                    report_id=report_id,
                    efforts=efforts,
                    keyword_efforts=keyword_efforts,
                )
                return {"message": f"Report {report_id} updated successfully"}

            # --- Meeting logic preserved ---
            meeting_message = "No meeting efforts processed"
            if user_id:
                meetings = await self.repository.fetch_meetings_by_user_id(user_id)
                if meetings:
                    saved_count = 0
                    for meeting in meetings:
                        cal_id = meeting.get("cal_id")
                        if not cal_id or await self.repository.meeting_id_exists(
                            cal_id
                        ):
                            continue
                        duration_minutes = meeting.get("duration_minutes", 0)
                        word_count = meeting.get("word_count", 0)
                        keywords_found = meeting.get("keyword", "")
                        keyword_count = meeting.get("repeated_keyword", 0)
                        await self.repository.save_meeting_report_data(
                            user_id=user_id,
                            org_id=org_id,
                            cal_id=cal_id,
                            word_count=word_count,
                            keywords_found=keywords_found,
                            keyword_count=keyword_count,
                            meeting_duration=duration_minutes,
                            efforts_time=duration_minutes,
                            created_by=user_id,
                        )
                        saved_count += 1
                    meeting_message = f"{saved_count} new meeting effort(s) stored for user_id {user_id}"
                else:
                    meeting_message = f"No meetings found for user_id {user_id}"

            # --- Validate user ---
            if not user_id or user_id <= 0:
                return {"message": "Invalid user_id. It must be a positive integer."}

            # --- Load rules and keywords ---
            rules = await self.repository.get_rules(org_id=org_id)
            if not rules:
                return {"message": f"No rules found for org_id {org_id}"}

            keyword_records = await self.repository.get_keywords(org_id)

            # --- Category helper ---
            async def get_category_for_keywords(
                keywords_found_unique: list[str],
                keyword_records: list[dict],
                categories: list[dict],
            ) -> int | None:
                """
                Determine the category ID for a set of keywords based on category priority.
                Now selects the category with the **lowest priority number**.
                """
                if not keywords_found_unique:
                    return None

                # Map keyword -> cat_id
                keyword_to_cat = {
                    kw["keyword_name"].strip().lower(): kw["cat_id"]
                    for kw in keyword_records
                    if kw.get("cat_id")
                }

                # Map cat_id -> priority
                cat_priority_map = {
                    cat["cat_id"]: cat.get("priority", float("inf"))
                    for cat in categories
                }

                # Collect categories found
                categories_found = set()
                for kw in keywords_found_unique:
                    cat_id = keyword_to_cat.get(kw.lower())
                    if cat_id is not None:
                        categories_found.add(cat_id)

                if not categories_found:
                    return None

                # Pick category with **lowest priority number**
                highest_priority_cat = None
                min_priority = float("inf")
                for cat_id in categories_found:
                    priority = cat_priority_map.get(cat_id, float("inf"))
                    if priority < min_priority:
                        min_priority = priority
                        highest_priority_cat = cat_id

                return highest_priority_cat

            # --- Keyword extraction helper (fixed) ---
            def extract_keywords(
                text: str, keyword_records: List[Dict[str, Any]]
            ) -> Tuple[List[str], List[str]]:
                all_keywords = []
                unique_keywords = set()
                if not text:
                    return [], []
                for record in keyword_records:
                    kw = record.get("keyword_name", "").lower()
                    if not kw:
                        continue
                    # find all occurrences (case-insensitive)
                    found = re.findall(
                        rf"(?<!\w){re.escape(kw)}(?!\w)",  # matches keyword not preceded/followed by word char
                        text,
                        flags=re.IGNORECASE,
                    )
                    if found:
                        # extend with the raw matched strings (keeps occurrence counts)
                        all_keywords.extend(found)
                        unique_keywords.add(kw)
                return list(unique_keywords), all_keywords

            # --- Process single mail ---
            async def process_single_mail(mail_data):
                mail_id = mail_data.get("mail_dtl_id")
                if not mail_id or await self.repository.mail_id_exists(mail_id):
                    return

                # BODY
                body_text = self._prepare_body_for_count_and_keywords(
                    mail_data.get("body") or ""
                )
                body_keywords_found_unique, body_all_keywords = extract_keywords(
                    body_text, keyword_records
                )
                body_word_count = (
                    len(re.findall(r"\b\w+\b", body_text)) if body_text else 0
                )
                body_effort = self.calculate_effort_dynamic(
                    body_word_count, body_all_keywords, is_attachment=False, rules=rules
                )

                # ATTACHMENTS - combine all attachments into a single text
                combined_attachment_text = ""
                attachments = await self.repository.get_attachments_by_mail_id(mail_id)
                if attachments:
                    for att in attachments:
                        text = await self._extract_attachment_text(att)
                        if text:
                            att_text_clean = self._prepare_body_for_count_and_keywords(
                                text
                            )
                            combined_attachment_text += att_text_clean + " "

                # Calculate combined attachment effort
                attach_word_count = (
                    len(re.findall(r"\b\w+\b", combined_attachment_text))
                    if combined_attachment_text
                    else 0
                )
                attach_found_unique, attach_all_keywords = extract_keywords(
                    combined_attachment_text, keyword_records
                )
                attach_best_effort = self.calculate_effort_dynamic(
                    attach_word_count,
                    attach_all_keywords,
                    is_attachment=True,
                    rules=rules,
                )
                attach_best_word_count = attach_word_count
                attach_best_all_keywords = (
                    attach_all_keywords  # <-- define it for later use
                )

                # Determine winner (body vs best attachment)
                if body_effort >= attach_best_effort:
                    actual_effort_time = body_effort
                    winning_all_keywords = body_all_keywords
                    winning_unique_keywords = body_keywords_found_unique
                    winning_word_count =attach_best_word_count
                else:
                    actual_effort_time = attach_best_effort
                    winning_all_keywords = attach_best_all_keywords
                    winning_unique_keywords = attach_found_unique  # you can use attach_found_unique or list(set(attach_best_all_keywords))
                    winning_word_count = attach_best_word_count  # <-- attachment wins, store its word count

                # keyword counts computed from winning_all_keywords (occurrences)
                total_keyword_count = len(winning_all_keywords)
                repeated_keyword_count = sum(
                    (v - 1)
                    for v in Counter([k.lower() for k in winning_all_keywords]).values()
                    if v > 1
                )

                planned_effort_time = max(
                    actual_effort_time, rules.get("minimum_effort")
                )

                categories = await self.repository.get_all_active_categories(org_id)
                category_id = await get_category_for_keywords(
                    winning_unique_keywords, keyword_records, categories
                )

                await self.repository.save_report_data(
                    user_id=user_id,
                    org_id=org_id,
                    mail_dtl_id=mail_id,
                    word_count=winning_word_count,
                    keywords_found=",".join(winning_unique_keywords),
                    keyword_count=total_keyword_count,
                    repeated_keyword_count=repeated_keyword_count,
                    actual_effort_time=actual_effort_time,
                    planned_effort_time=planned_effort_time,
                    cat_id=category_id,
                    created_by=mail_data.get("created_by"),
                )

            # --- Single mail ---
            if mail_dtl_id:
                mail_rows = await self.repository.fetch_mail_details_by_id(mail_dtl_id)
                if not mail_rows:
                    return {"message": f"No email found for mail_dtl_id {mail_dtl_id}"}
                await process_single_mail(mail_rows[0])
                return {"message": "Effort stored successfully"}

            # --- Bulk mail ---
            mail_data_list = await self.repository.get_mail_details(user_id)
            if not mail_data_list:
                return {"message": f"No emails found for user_id {user_id}"}
            existing_mail_ids = await self.repository.get_existing_mail_ids(
                [m.get("mail_dtl_id") for m in mail_data_list if m.get("mail_dtl_id")]
            )

            for mail_data in mail_data_list:
                if (
                    not mail_data.get("mail_dtl_id")
                    or mail_data.get("mail_dtl_id") in existing_mail_ids
                ):
                    continue
                await process_single_mail(mail_data)

            return {"message": "Effort stored successfully"}

        except Exception as e:
            logger.error(f"Error processing mail effort: {str(e)}")
            return {"message": f"Error processing mail effort: {str(e)}"}

    # --- Effort calculation helper ---
    def calculate_effort_dynamic(
        self,
        word_count: int,
        all_keywords: List[str],
        is_attachment: bool,
        rules: Dict[str, Any],
    ) -> float:
        words_per_unit = (
            rules["attachment_word_count"]
            if is_attachment
            else rules["email_body_word_count"]
        )
        minutes_per_unit = (
            rules["attachment_effort"] if is_attachment else rules["email_body_effort"]
        )
        keyword_repeat_effort = rules["keyword_repeat_effort"]
        minimum_effort = rules["minimum_effort"]

        base_minutes = (word_count * minutes_per_unit) / words_per_unit

        repeated_count = sum(
            (v - 1)
            for v in Counter([kw.lower() for kw in all_keywords]).values()
            if v > 1
        )
        bonus_minutes = repeated_count * keyword_repeat_effort

        total_effort = base_minutes + bonus_minutes
        return max(total_effort, minimum_effort)

    async def _extract_attachment_text(self, attachment: Dict[str, Any]) -> str:
        """
        Extract text from an attachment file of any common type.
        Supports PDF, DOCX, PPTX, CSV, TXT, XLSX, and tries fallback for unknown types.
        """
        try:
            file_type = attachment.get("attach_type")
            file_content = attachment.get("content")  # bytes
            if not file_content:
                return ""

            text = ""

            # PDF
            if file_type == "application/pdf":
                reader = PyPDF2.PdfReader(io.BytesIO(file_content))
                for page in reader.pages:
                    page_text = page.extract_text()
                    if page_text:
                        text += page_text + " "

            # Word DOCX
            elif (
                file_type
                == "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
            ):
                doc = Document(io.BytesIO(file_content))
                for para in doc.paragraphs:
                    text += para.text + " "

            # PowerPoint PPTX
            elif (
                file_type
                == "application/vnd.openxmlformats-officedocument.presentationml.presentation"
            ):
                ppt = Presentation(io.BytesIO(file_content))
                for slide in ppt.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text += shape.text + " "

            # XLSX
            elif (
                file_type
                == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
            ):
                wb = load_workbook(filename=io.BytesIO(file_content), data_only=True)
                for sheet in wb.worksheets:
                    for row in sheet.iter_rows(values_only=True):
                        row_text = " ".join(
                            [str(cell) for cell in row if cell is not None]
                        )
                        text += row_text + " "
            
            # XLS (Excel 97-2003)
            elif file_type == "application/vnd.ms-excel":
                wb = xlrd.open_workbook(file_contents=file_content)
                for sheet in wb.sheets():
                    for row_idx in range(sheet.nrows):
                        row_values = [str(sheet.cell_value(row_idx, col)) for col in range(sheet.ncols)]
                        row_text = " ".join(row_values)
                        text += row_text + " "

            # CSV
            elif file_type in ("text/csv", "application/csv"):
                decoded = (
                    file_content.decode("utf-8", errors="ignore")
                    if isinstance(file_content, bytes)
                    else file_content
                )
                f = io.StringIO(decoded)
                csv_reader = csv.reader(f)
                for row in csv_reader:
                    # Combine cells with a space, ignore empty cells
                    row_text = " ".join([str(cell) for cell in row if cell])
                    text += row_text + " "

            # Plain text or other text/*
            elif file_type.startswith("text/"):
                text = file_content.decode("utf-8", errors="ignore")

            # Fallback: try UTF-8 decoding for unknown files
            else:
                try:
                    text = file_content.decode("utf-8", errors="ignore")
                except Exception:
                    text = ""

            # Clean multiple spaces
            text = " ".join(text.split())
            return text.strip()

        except Exception as e:
            logger.error(f"Error extracting attachment text: {str(e)}")
            return ""

    def _prepare_body_for_count_and_keywords(self, body: str) -> str:
        if not body:
            return ""
        cleaned = re.sub(
            r"\s+", " ", body.replace("\n", " ").replace("\t", " ").strip()
        )
        return cleaned

    def _extract_keywords(
        self, text: str, keyword_records: List[Dict[str, Any]]
    ) -> Tuple[List[str], List[str]]:
        combined_text = text.lower()
        found_keywords_unique = []
        all_keywords = []

        for kw in keyword_records:
            kw_text = kw["keyword_name"].lower().strip()
            if not kw_text:
                continue
            # Use negative lookbehind/lookahead instead of \b to handle special chars
            matches = re.findall(rf"(?<!\w){re.escape(kw_text)}(?!\w)", combined_text)
            if matches:
                found_keywords_unique.append(kw_text)
                all_keywords.extend([kw_text] * len(matches))

        return found_keywords_unique, all_keywords

    # def _calculate_effort(
    #     self, body_text: str, keywords: list, rules: dict, is_attachment: bool = False
    # ) -> float:
    #     """
    #     Calculate effort dynamically based on rules from DB.
    #     All values are fetched from `rules` dict; no hardcoded defaults.
    #     """

    #     # Determine which keys to use
    #     if is_attachment:
    #         words_per_unit = rules["attachment_word_count"]
    #         minutes_per_unit = rules["attachment_effort"]
    #     else:
    #         words_per_unit = rules["email_body_word_count"]
    #         minutes_per_unit = rules["email_body_effort"]

    #     minutes_per_keyword = rules["keyword_repeat_effort"]
    #     minimum_effort = rules["minimum_effort"]

    #     # Prepare text and split into words
    #     clean_body = self._prepare_body_for_count_and_keywords(body_text)
    #     word_list = [w.lower().strip(",.!?") for w in clean_body.split()]

    #     # Base effort by word count
    #     base_minutes = (len(word_list) / words_per_unit) * minutes_per_unit

    #     # Bonus for repeated keywords
    #     bonus_minutes = 0
    #     word_counter = Counter(word_list)
    #     for kw in keywords:
    #         count = word_counter[kw.lower()]
    #         if count > 1:
    #             bonus_minutes += (count - 1) * minutes_per_keyword

    #     total_effort = base_minutes + bonus_minutes

    #     # Apply minimum effort
    #     if total_effort < minimum_effort:
    #         total_effort = minimum_effort

    #     return total_effort

    async def _extract_pdf_text(self, pdf_files: List[Dict]) -> List[str]:
        """
        Extract text from PDF files.
        Returns a list of extracted text strings.
        """
        extracted_texts = []

        for pdf_file in pdf_files:
            try:
                file_content = pdf_file.get("content")
                if not file_content:
                    logger.warning("PDF file has no content")
                    continue

                pdf_stream = io.BytesIO(file_content)
                pdf_reader = PyPDF2.PdfReader(pdf_stream)
                text = ""

                if len(pdf_reader.pages) == 0:
                    logger.warning("PDF file has no pages")
                    continue

                for i in range(len(pdf_reader.pages)):
                    try:
                        page = pdf_reader.pages[i]
                        page_text = page.extract_text()
                        if page_text:
                            text += page_text + " "
                    except (IndexError, Exception) as page_error:
                        logger.error(
                            f"Error extracting text from PDF page {i}: {str(page_error)}"
                        )
                        continue

                if text.strip():
                    extracted_texts.append(text.strip())
                    logger.info(
                        f"Successfully extracted {len(text.strip())} characters from PDF"
                    )

            except Exception as e:
                logger.error(f"Error extracting text from PDF: {str(e)}")
                continue

        return extracted_texts

    async def _extract_csv_text(self, csv_files: List[Dict]) -> List[str]:
        """
        Extract text from CSV files.
        Returns a list of extracted text strings.
        """
        extracted_texts = []

        for csv_file in csv_files:
            try:
                file_content = csv_file.get("content")
                if not file_content:
                    logger.warning("CSV file has no content")
                    continue

                if isinstance(file_content, bytes):
                    file_content = file_content.decode("utf-8")

                csv_stream = io.StringIO(file_content)
                csv_reader = csv.reader(csv_stream)
                text_parts = []

                for row in csv_reader:
                    if row:
                        text_parts.extend(row)

                if text_parts:
                    extracted_texts.append(" ".join(text_parts))
                    logger.info(
                        f"Successfully extracted {len(text_parts)} text parts from CSV"
                    )

            except Exception as e:
                logger.error(f"Error extracting text from CSV: {str(e)}")
                continue

        return extracted_texts

    async def fetch_mail_report_data(
        self,
        from_date: Union[date, str, None] = None,
        to_date: Union[date, str, None] = None,
        user_id: Union[int, None] = None,
        role_id: Union[int, None] = None,
        org_id: Union[int, None] = None,
        keyword: Optional[str] = None,
        category: Optional[str] = None,
        user: Optional[str] = None,
        file_type: Optional[str] = None,
    ):
        try:
            filters = {
                "keyword": keyword,
                "category": category,
                "user": user,
                "file_type": file_type,
            }

            if user_id and user_id > 0:
                if role_id == 1:  # Admin
                    results, total_count = (
                        await self.repository.get_report_data_by_date_admin(
                            from_date, to_date, user_id, org_id, filters
                        )
                    )
                else:
                    results, total_count = (
                        await self.repository.get_report_data_by_date_userId(
                            from_date, to_date, user_id, filters
                        )
                    )

            else:
                results, total_count = await self.repository.get_report_data_by_date(
                    from_date, to_date, filters
                )

            return {
                "success": True,
                "data": results,
                "count": total_count,
                "from_date": from_date.isoformat(),
                "to_date": to_date.isoformat(),
            }
        except Exception as e:
            logger.error(f"Error fetching report data: {str(e)}")
            return {"success": False, "message": str(e)}

    async def fetch_mail_details_by_id(self, mailDtl_id: int) -> Dict[str, Any]:
        """
        Service to fetch mail details by id.
        Returns subject, body, and list of attachment names.
        """
        rows = await self.repository.fetch_mail_details_by_id(mailDtl_id)

        if not rows:
            return {"message": "No mail details found", "data": None}

        user_id = rows[0]["user_id"]
        mailDtl_id = rows[0]["mail_dtl_id"]
        subject = rows[0]["subject"]
        body = rows[0]["body"]
        attachments = [row["attach_name"] for row in rows if row["attach_name"]]

        return {
            "user_id": user_id,
            "mailDtl_id": mailDtl_id,
            "subject": subject,
            "body": body,
            "attachments": attachments,
        }

    # ---------------------- Meeting Details(preview)---------------------
    async def fetch_meeting_details_by_id(self, cal_id: int) -> Dict[str, Any]:
        """
        Service to fetch meeting details by id.
        Returns subject, body, and list of attachment names.
        """
        rows = await self.repository.fetch_meeting_details_by_id(cal_id)

        if not rows:
            return {"message": "No meeting details found", "data": None}

        user_id = rows[0]["user_id"]
        cal_id = rows[0]["cal_id"]
        subject = rows[0]["title"]
        body = rows[0]["description"]
        attachments = [row["attach_name"] for row in rows if row["attach_name"]]

        return {
            "user_id": user_id,
            "cal_id": cal_id,
            "subject": subject,
            "body": body,
            "attachments": attachments,
        }

    async def preview_attachment(
        self, mail_dtl_id: int, user_id: int, attachment_name: str
    ):
        record = await self.repository.get_attachment(
            mail_dtl_id, user_id, attachment_name
        )
        if not record:
            return None

        file_path = record["attach_path"]
        if not os.path.exists(file_path):
            return None

        return {
            "path": file_path,
            "type": record["attach_type"],
            "name": record["attach_name"],
        }

    async def insert_rule(self, org_id: int, rule_key: str, rule_value: float):
        try:
            msg = await self.repository.insert_rule(org_id, rule_key, rule_value)
            return msg
        except Exception as e:
            logger.error(f"Error inserting rule: {e}")
            return {"message": f"Error inserting rule: {str(e)}"}

    async def fetch_business_rules(self, org_id: int, user_id: int) -> List[Dict]:
        """
        Business logic to fetch multiple business rules for org.
        """
        rules = await self.repository.get_business_rules_by_org_id(org_id)

        if not rules:
            return []

        return [
            {
                "rule_id": rule["rule_id"],
                "rule_name": rule["rule_key"],
                "rule_value": rule["rule_value"],
            }
            for rule in rules
        ]

    async def update_rule(self, rule_id: int, value: float, org_id: int) -> bool:
        # Validate inputs again if needed
        if value <= 0 or rule_id <= 0:
            return False

        updated_rows = await self.repository.update_rule_value(rule_id, value, org_id)
        return updated_rows > 0

    # --------------------- Fetch Meeting Data ---------------------
    async def fetch_meeting_data(
        self,
        org_id: int,
        from_date: Union[date, str, None] = None,
        to_date: Union[date, str, None] = None,
        user_id: int | None = None,
        page: int = None,
        limit: int = None,
        role_id: int | None = None,
        user: Optional[str] = None,
    ):
        try:
            if not role_id:
                return {"success": False, "message": "role_id is required"}

            offset = (page - 1) * limit

            if role_id == 1:  # Admin
                results, total_count = await self.repository.fetch_meeting_data(
                    org_id,
                    from_date,
                    to_date,
                    None,
                    limit,
                    offset,
                    is_admin=True,
                    user=user,
                )
            else:
                if not user_id:
                    return {"success": False, "message": "user_id is required"}
                results, total_count = await self.repository.fetch_meeting_data(
                    org_id,
                    from_date,
                    to_date,
                    user_id,
                    limit,
                    offset,
                    is_admin=False,
                    user=None,
                )

            return {
                "success": True,
                "data": results,
                "total": total_count,
                "page": page,
                "limit": limit,
                "totalPages": (total_count + limit - 1) // limit,  # ceil division
                "from_date": str(from_date) if from_date else None,
                "to_date": str(to_date) if to_date else None,
            }
        except Exception as e:
            logger.error(f"Error fetching meeting data: {str(e)}")
            return {"success": False, "message": str(e)}

    # -------------------- Fetch Entity Types --------------------
    async def fetch_entity_types(self):
        try:
            entities = await self.repository.get_all_entity_types()
            return {"success": True, "data": entities}
        except Exception as e:
            return {"success": False, "message": str(e)}

    # -------------------- Update Meeting Effort --------------------
    async def update_meeting_effort(
        self, meeting_report_id: int, efforts: Optional[float]
    ):
        if efforts is not None and efforts < 0:
            raise HTTPException(
                status_code=400, detail="Effort time cannot be negative"
            )

        await self.repository.update_meeting_effort(meeting_report_id, efforts)
        return {"message": "Meeting effort updated successfully"}

    # ---------------------fetch folder name------------------
    async def fetch_folders_name(self) -> List[str]:
        return await self.repository.fetch_folders_name()

    # ---------------recalculate the effort time with new business rule-----------------
    async def recalc_effort_for_existing_reports(
        self, org_id: int, user_id: Optional[int] = None
    ):
        """
        Recalculate effort for all previously synced emails (reports) for the given org_id
        and optionally filtered by user_id. Only updates actual_effort_time and planned_effort_time.
        """
        try:
            # Load the latest business rules
            rules = await self.repository.get_rules(org_id=org_id)
            if not rules:
                return {"message": f"No rules found for org_id {org_id}"}

            # Load all keywords for this org
            keyword_records = await self.repository.get_keywords(org_id)

            # Fetch all reports to recalc
            reports = await self.repository.get_all_reports(
                org_id
            )  # implement this repository method
            if not reports:
                return {"message": "No reports found to recalculate"}

            for report_row in reports:
                mail_id = report_row.get("mail_dtl_id")
                if not mail_id:
                    continue

                # --- BODY processing ---
                body_text = self._prepare_body_for_count_and_keywords(
                    report_row.get("body") or ""
                )
                body_keywords_found_unique, body_all_keywords = self._extract_keywords(
                    body_text, keyword_records
                )
                body_word_count = (
                    len(re.findall(r"\b\w+\b", body_text)) if body_text else 0
                )
                body_effort = self.calculate_effort_dynamic(
                    body_word_count, body_all_keywords, is_attachment=False, rules=rules
                )

                # --- ATTACHMENTS processing ---
                combined_attachment_text = ""
                attachments = await self.repository.get_attachments_by_mail_id(mail_id)
                if attachments:
                    for att in attachments:
                        att_text = await self._extract_attachment_text(att)
                        if att_text:
                            att_text_clean = self._prepare_body_for_count_and_keywords(
                                att_text
                            )
                            combined_attachment_text += att_text_clean + " "

                attach_word_count = (
                    len(re.findall(r"\b\w+\b", combined_attachment_text))
                    if combined_attachment_text
                    else 0
                )
                attach_found_unique, attach_all_keywords = self._extract_keywords(
                    combined_attachment_text, keyword_records
                )
                attach_effort = self.calculate_effort_dynamic(
                    attach_word_count,
                    attach_all_keywords,
                    is_attachment=True,
                    rules=rules,
                )

                # --- Determine final effort ---
                if body_effort >= attach_effort:
                    actual_effort_time = body_effort
                else:
                    actual_effort_time = attach_effort

                planned_effort_time = max(
                    actual_effort_time, rules.get("minimum_effort")
                )

                # --- Update report with new effort only ---
                await self.repository.update_report_efforts_only(
                    report_id=report_row["report_id"],
                    actual_effort_time=actual_effort_time,
                    planned_effort_time=planned_effort_time,
                )

            return {
                "message": f"Recalculated efforts for {len(reports)} reports successfully"
            }

        except Exception as e:
            logger.error(f"Error recalculating report efforts: {str(e)}")
            return {"message": f"Error recalculating report efforts: {str(e)}"}

    # --------------------- Fetch Keywords List --------------------
    async def fetch_keywords_list(self, org_id: int) -> List[str]:
        if org_id <= 0:
            raise ValueError("org_id must be a positive integer")
        try:
            keywords = await self.repository.fetch_keywords_list(org_id)
            return keywords
        except Exception as e:
            logger.error(f"Error fetching keywords: {str(e)}")
            return {"success": False, "message": str(e)}

    # --------------------- Fetch Category List --------------------
    async def fetch_category_list(self, org_id: int) -> List[str]:
        if org_id <= 0:
            raise ValueError("org_id must be a positive integer")
        try:
            categories = await self.repository.fetch_category_list(org_id)
            return categories
        except Exception as e:
            logger.error(f"Error fetching categories: {str(e)}")
            return {"success": False, "message": str(e)}

    # --------------------- Fetch Users List -----------------------
    async def fetch_users_list(self, org_id: int, role_id:int) -> List[str]:
        if org_id <= 0:
            raise ValueError("org_id must be a positive integer")
        try:
            users = await self.repository.fetch_users_list(org_id, role_id)
            return users
        except Exception as e:
            logger.error(f"Error fetching users: {str(e)}")
            return {"success": False, "message": str(e)}

    # --------------------- Fetch File Type List -------------------
    async def fetch_fileType_list(self) -> List[str]:
        try:
            file_types = await self.repository.fetch_fileType_list()
            return file_types
        except Exception as e:
            logger.error(f"Error fetching file types: {str(e)}")
            return {"success": False, "message": str(e)}
